import io
import os
import re
import html
import gc
import time
import csv
import json
import math
import string
import zipfile
import tempfile
import logging
import secrets
from dataclasses import dataclass, field
from urllib.parse import urlparse
from collections import Counter, defaultdict
from typing import Dict, List, Tuple, Iterable, Optional, Callable, Any, Union, Set
from datetime import datetime

import numpy as np
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from io import BytesIO
from wordcloud import WordCloud, STOPWORDS
from matplotlib import font_manager
from itertools import pairwise
import openai

# --- graph imports
import networkx as nx
import networkx.algorithms.community as nx_comm
from streamlit_agraph import agraph, Node, Edge, Config

# -3rd party imports checks
try:
    import requests
    from bs4 import BeautifulSoup
except ImportError:
    requests = None
    BeautifulSoup = None

try:
    import qrcode
except ImportError:
    qrcode = None

try:
    from scipy.stats import beta as beta_dist
except ImportError:
    beta_dist = None

try:
    from sklearn.decomposition import LatentDirichletAllocation, NMF
    from sklearn.feature_extraction import DictVectorizer
except ImportError:
    LatentDirichletAllocation = None
    NMF = None
    DictVectorizer = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from dateutil import parser as date_parser
except ImportError:
    date_parser = None

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import nltk
    from nltk.sentiment.vader import SentimentIntensityAnalyzer
    from nltk.stem import WordNetLemmatizer
except ImportError:
    nltk = None
    SentimentIntensityAnalyzer = None
    WordNetLemmatizer = None


# ⚙️ constants and config
# ========================================

MAX_TOPIC_DOCS = 50_000
MAX_SPEAKER_NAME_LENGTH = 30
SENTIMENT_ANALYSIS_TOP_N = 5000
URL_SCRAPE_RATE_LIMIT_SECONDS = 1.0
PROGRESS_UPDATE_MIN_INTERVAL = 100
NPMI_MIN_FREQ = 3
MAX_FILE_SIZE_MB = 200

# regex patterns
HTML_TAG_RE = re.compile(r"<[^>]+>")
CHAT_ARTIFACT_RE = re.compile(
    r":\w+:"
    r"|\b(?:monday|tuesday|wednesday|thursday|friday|saturday|sunday|today|yesterday) at \d{1,2}:\d{2}\b"
    r"|\b\d+\s+repl(?:y|ies)\b"
    r"|\d{2}:\d{2}:\d{2}\.\d{3}\s+-->\s+\d{2}:\d{2}:\d{2}\.\d{3}"
    r"|\[[^\]]+\]",
    flags=re.IGNORECASE
)
URL_EMAIL_RE = re.compile(
    r'(?:https?://|www\.)[a-zA-Z0-9-]+(?:\.[a-zA-Z0-9-]+)+[^\s]*'
    r'|(?:[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})',
    flags=re.IGNORECASE
)
# NER Pattern: matches 3 types of entities:
# acronyms (2+ Uppercase): DARPA, EU, NATO
# complex IDs (Cap start + digits/hyphens): COVID-19, Mi-6, G-7
# standard proper nouns (title case phrases): John Doe, Project Gutenberg
NER_CAPS_RE = re.compile(
    r'\b(?:'
    r'[A-Z]{2,}'                         # acronyms (DARPA)
    r'|[A-Z][a-zA-Z0-9-]*\d[a-zA-Z0-9-]*' # complex IDs (COVID-19)
    r'|[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*'    # standard names (John Doe)
    r')\b'
)

# logger setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("IntelEngine")

# custom exceptions
class ReaderError(Exception):
    pass

class ValidationError(Exception):
    pass

# 📦 dataclassed
# ===========================================

@dataclass
class CleaningConfig:
    remove_chat: bool = True
    remove_html: bool = True
    remove_urls: bool = True
    unescape: bool = True
    phrase_pattern: Optional[re.Pattern] = None

@dataclass
class ProcessingConfig:
    min_word_len: int = 2
    drop_integers: bool = True
    compute_bigrams: bool = True
    use_lemmatization: bool = False
    translate_map: Dict[int, Optional[int]] = field(default_factory=dict)
    stopwords: Set[str] = field(default_factory=set)


# 🛡️ security and validation utils
# ==========================================

def get_auth_password() -> str:
    pwd = st.secrets.get("auth_password")
    if not pwd:
        st.error("🚨 Configuration Error: 'auth_password' not set in .streamlit/secrets.toml.")
        st.stop()
    return pwd

def validate_url(url: str) -> bool:
    try:
        parsed = urlparse(url)
        if parsed.scheme not in ('http', 'https'):
            return False
        if parsed.hostname in ('localhost', '127.0.0.1', '0.0.0.0', '::1'):
            return False
        return True
    except Exception:
        return False

def validate_sketch_data(data: Dict) -> bool:
    REQUIRED_KEYS = {"total_rows", "counts", "bigrams", "topic_docs"}
    if not isinstance(data, dict): return False
    if not REQUIRED_KEYS.issubset(data.keys()): return False
    return True


# 🧠 core logic, scanner
# =========================================

class StreamScanner:
    def __init__(self, doc_batch_size=5):
        self.global_counts = Counter()
        self.global_bigrams = Counter()
        self.total_rows_processed = 0
        self.topic_docs: List[Counter] = []
        
        # newer analytical structures (v2.0+)
        self.temporal_counts = defaultdict(Counter) # { '2023-10-01': {'word': 5} }
        self.category_counts = defaultdict(Counter) # { 'CategoryA': {'word': 10} }
        self.doc_freqs = Counter() # DF for TF-IDF
        self.entity_counts = Counter() # NER lite storage
        
        self.DOC_BATCH_SIZE = doc_batch_size
        self.limit_reached = False

    def set_batch_size(self, size: int):
        self.DOC_BATCH_SIZE = size

    def update_global_stats(self, counts: Counter, bigrams: Counter, rows: int):
        self.global_counts.update(counts)
        self.global_bigrams.update(bigrams)
        self.total_rows_processed += rows
    
    # only updates doc_freqs if still sampling (prevents unnecessary work)    
    def add_topic_sample(self, doc_counts: Counter):
        if not doc_counts: return
        
        if not self.limit_reached and len(self.topic_docs) < MAX_TOPIC_DOCS:
            self.doc_freqs.update(doc_counts.keys())
            self.topic_docs.append(doc_counts)
            if len(self.topic_docs) >= MAX_TOPIC_DOCS:
                self.limit_reached = True
        # if limit already reached, do nothing — saves memory/CPU
       
    def update_metadata_stats(self, date_key: Optional[str], cat_key: Optional[str], tokens: List[str]):
        if date_key:
            self.temporal_counts[date_key].update(tokens)
        if cat_key:
            self.category_counts[cat_key].update(tokens)

    def update_entities(self, entities: List[str]):
        if entities:
            self.entity_counts.update(entities)

    def to_json(self) -> str:
        # simplifying complex structures for JSON serialization
        serializable_bigrams = {f"{k[0]}|{k[1]}": v for k, v in self.global_bigrams.items()}
        data = {
            "total_rows": self.total_rows_processed,
            "counts": dict(self.global_counts),
            "bigrams": serializable_bigrams,
            "topic_docs": [dict(c) for c in self.topic_docs],
            "limit_reached": self.limit_reached,
            # Persistence for new features
            "temporal_counts": {k: dict(v) for k, v in self.temporal_counts.items()},
            "entity_counts": dict(self.entity_counts),
            "doc_freqs": dict(self.doc_freqs)
        }
        return json.dumps(data)

    def load_from_json(self, json_str: str) -> bool:
        try:
            data = json.loads(json_str)
            if not validate_sketch_data(data): return False
            
            self.total_rows_processed = data.get("total_rows", 0)
            self.global_counts = Counter(data.get("counts", {}))
            
            raw_bigrams = data.get("bigrams", {})
            self.global_bigrams = Counter()
            for k, v in raw_bigrams.items():
                if "|" in k:
                    parts = k.split("|", 1)
                    self.global_bigrams[(parts[0], parts[1])] = v
            
            self.topic_docs = [Counter(d) for d in data.get("topic_docs", [])]
            self.limit_reached = data.get("limit_reached", False)
            
            # new
            self.entity_counts = Counter(data.get("entity_counts", {}))
            self.doc_freqs = Counter(data.get("doc_freqs", {}))
            
            raw_temp = data.get("temporal_counts", {})
            self.temporal_counts = defaultdict(Counter)
            for k, v in raw_temp.items():
                self.temporal_counts[k] = Counter(v)
                
            return True
        except Exception as e:
            logger.error(f"JSON Load Error: {e}")
            return False

# 📈 maturity modeling logic (Multi-persona)
# ==========================================

# 📈 maturity modeling logic (multi-Persona)
# ==========================================

class MaturityAssessor:
    """
    Evaluates maturity based on linguistic markers using selectable 'Personas'.
    Supports switching between Business, EdTech, and Policy contexts.

    Now supports:
    - Single-word tokens (unigrams)
    - Multi-word phrases (bigrams), via the scanner's global_bigrams
    """
    def __init__(self):
        # library of models
        # NOTE: each level can now optionally define:
        #   - "terms": set[str]          -> single tokens
        #   - "phrases": list[str]       -> multi-word phrases (space-separated)
        self.models = {
            "🏫 EdTech & LMS Ops": {
                "desc": "Evaluates LMS utilization from 'Digital Repository' (L1) to 'Connected Ecosystem' (L5).",
                "levels": {
                    1: {
                        "name": "Digital Repository (Static)",
                        "color": "#d62728",
                        "terms": {
                            "upload", "download", "pdf", "file", "syllabus", "login", "password",
                            "access", "content", "link", "ppt", "doc", "email", "submit", "paper", "static"
                        },
                        "phrases": [
                            "course shell", "file repository", "content dump"
                        ]
                    },
                    2: {
                        "name": "Managed Courseware (Tools)",
                        "color": "#ff7f0e",
                        "terms": {
                            "quiz", "gradebook", "discussion", "rubric", "module", "assignment",
                            "calendar", "announcement", "feedback", "group", "template", "checklist", "forum"
                        },
                        "phrases": [
                            "online quiz", "assignment submission", "discussion forum"
                        ]
                    },
                    3: {
                        "name": "Integrated (Connected)",
                        "color": "#f7b731",
                        "terms": {
                            "lti", "integration", "api", "plugin", "interoperability", "tool",
                            "external", "sso", "vendor", "connect", "ecosystem", "zoom", "teams", "turnitin", "scorm"
                        },
                        "phrases": [
                            "learning tools", "third party", "external tool", "deep integration"
                        ]
                    },
                    4: {
                        "name": "Data-Informed (Adaptive)",
                        "color": "#2ca02c",
                        "terms": {
                            "analytics", "engagement", "retention", "risk", "dashboard", "report",
                            "outcome", "competency", "mastery", "release", "adaptive", "personalized", "pathway", "agent"
                        },
                        "phrases": [
                            "learning analytics", "course analytics", "early alert", "release conditions"
                        ]
                    },
                    5: {
                        "name": "Optimized (Strategic)",
                        "color": "#9467bd",
                        "terms": {
                            "governance", "strategy", "accessibility", "udl", "equity", "inclusion",
                            "continuous", "scale", "innovation", "agency", "holistic", "success", "lifelong", "transform"
                        },
                        "phrases": [
                            "governance framework", "continuous improvement", "student success", "institutional strategy"
                        ]
                    }
                }
            },
            "🏢 General Business Ops": {
                "desc": "Standard CMMI model: From 'Ad-Hoc' chaos to 'Optimized' strategy.",
                "levels": {
                    1: {
                        "name": "Ad-Hoc / Reactive",
                        "color": "#d62728",
                        "terms": {
                            "urgent", "fix", "panic", "broken", "late", "fail", "incident", "manual", "fire", "chaos"
                        },
                        "phrases": [
                            "fire drill", "last minute", "workaround"
                        ]
                    },
                    2: {
                        "name": "Managed / Project",
                        "color": "#ff7f0e",
                        "terms": {
                            "plan", "track", "project", "deadline", "schedule", "assign", "meeting", "status", "budget"
                        },
                        "phrases": [
                            "project plan", "status report", "project charter"
                        ]
                    },
                    3: {
                        "name": "Defined / Standardized",
                        "color": "#f7b731",
                        "terms": {
                            "standard", "process", "policy", "document", "compliance", "audit", "workflow", "consistent"
                        },
                        "phrases": [
                            "standard operating", "standardized process", "policy framework"
                        ]
                    },
                    4: {
                        "name": "Measured / Quantitative",
                        "color": "#2ca02c",
                        "terms": {
                            "metric", "kpi", "measure", "data", "analysis", "trend", "dashboard", "roi", "forecast"
                        },
                        "phrases": [
                            "key performance indicator", "data driven", "variance analysis"
                        ]
                    },
                    5: {
                        "name": "Optimizing / Strategic",
                        "color": "#9467bd",
                        "terms": {
                            "innovate", "strategy", "vision", "culture", "synergy", "scale", "optimize", "best-in-class"
                        },
                        "phrases": [
                            "strategic roadmap", "continuous improvement", "organizational transformation"
                        ]
                    }
                }
            },
            "⚖️ Policy & Governance": {
                "desc": "Evaluates policy maturity from 'Reactive/Enforcement' (L1) to 'Systemic/Holistic' (L5).",
                "levels": {
                    1: {
                        "name": "Enforcement / Reactive",
                        "color": "#d62728",
                        "terms": {
                            "violation", "sanction", "ban", "prohibit", "force", "threat", "danger",
                            "emergency", "crisis", "incident", "restriction", "penalty", "risk", "security"
                        },
                        "phrases": [
                            "zero tolerance", "strict enforcement"
                        ]
                    },
                    2: {
                        "name": "Procedural / Draft",
                        "color": "#ff7f0e",
                        "terms": {
                            "draft", "proposal", "clause", "article", "amendment", "review",
                            "committee", "meeting", "agenda", "timeline", "signature", "ratify", "consensus"
                        },
                        "phrases": [
                            "draft policy", "working group", "policy proposal"
                        ]
                    },
                    3: {
                        "name": "Operational / Implemented",
                        "color": "#f7b731",
                        "terms": {
                            "framework", "mechanism", "guideline", "standard", "monitor",
                            "verify", "compliance", "report", "mandate", "coordination", "treaty"
                        },
                        "phrases": [
                            "compliance framework", "implementation plan", "governance mechanism"
                        ]
                    },
                    4: {
                        "name": "Evidence-Based / Analysis",
                        "color": "#2ca02c",
                        "terms": {
                            "assessment", "evaluation", "impact", "data", "research", "finding",
                            "indicator", "measure", "trend", "forecast", "efficacy", "evidence"
                        },
                        "phrases": [
                            "impact assessment", "evidence based", "data driven"
                        ]
                    },
                    5: {
                        "name": "Systemic / Sustainable",
                        "color": "#9467bd",
                        "terms": {
                            "sustainable", "resilient", "holistic", "global", "ecosystem",
                            "peace", "development", "cooperation", "future", "inclusive", "norms", "universal"
                        },
                        "phrases": [
                            "holistic approach", "systemic change", "sustainable development"
                        ]
                    }
                }
            },
            "🎓 Brightspace Maturity Model (12-Domain)": {
    "desc": "Official 12-domain Brightspace Admin Maturity Model. Evaluates from Foundational to Leading Edge across Platform Admin, Curriculum, Student Engagement, Analytics, Assessment, Instructor Efficiency, Change Management, Knowledge Management, Accessibility, User Support, Innovation, and Collaboration.",
    "type": "domain_based",
    "domains": {
        "01_platform_admin": {
            "name": "Platform & Technical Administration",
            "short": "Platform Admin",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "upload", "download", "csv", "manual", "default", "password",
                        "login", "enrol", "enroll", "patching", "basic", "static",
                        "broad", "minimal", "limited"
                    },
                    "phrases": [
                        "course shell", "default settings", "csv uploads",
                        "manual updates", "default configurations",
                        "minimal customization", "basic password",
                        "manual user", "limited api", "limited use"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "api", "automation", "sso", "lti", "scorm", "xapi",
                        "scripting", "compliance", "dashboard", "provisioning",
                        "workflow", "documentation", "blueprint", "naming",
                        "monitoring", "wcag", "gdpr", "impersonation"
                    },
                    "phrases": [
                        "role based", "access control", "data hub",
                        "naming conventions", "least privilege",
                        "automated user", "course provisioning",
                        "standard operating", "regular review",
                        "documented processes", "sis integration"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "governance", "middleware", "forecasting", "transformation",
                        "scalable", "monitoring", "capacity", "branded",
                        "lifecycle", "archiving", "penetration", "siem"
                    },
                    "phrases": [
                        "continuous improvement", "capacity planning",
                        "custom scripts", "integration roadmap",
                        "real-time monitoring", "power bi", "tableau",
                        "automation strategy", "performance monitoring",
                        "full automation", "cross-platform automation",
                        "system health"
                    ]
                }
            }
        },
        "02_curriculum": {
            "name": "Curriculum Development & Delivery",
            "short": "Curriculum",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "syllabus", "outline", "upload", "organize",
                        "template", "schedule", "guideline", "accessible"
                    },
                    "phrases": [
                        "learning objectives", "course materials",
                        "curriculum goals", "basic accessibility",
                        "course design", "content effectively"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "instructional", "multimedia", "interactive",
                        "rubric", "inclusive", "formative", "differentiated",
                        "personalize"
                    },
                    "phrases": [
                        "instructional design", "active learning",
                        "diverse learning", "release conditions",
                        "learning styles", "course content"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "accreditation", "competency", "programmatic",
                        "innovation", "pedagogical"
                    },
                    "phrases": [
                        "program-level learning", "cross-functional teams",
                        "competency-based learning", "continuous improvement",
                        "curriculum evaluation", "institutional goals",
                        "curriculum improvements"
                    ]
                }
            }
        },
        "03_student_engagement": {
            "name": "Student Engagement & Success",
            "short": "Student Engagement",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "navigate", "announcement", "discussion", "tutorial",
                        "onboarding", "calendar", "login", "access"
                    },
                    "phrases": [
                        "course shell", "technical support",
                        "student learning", "access issues",
                        "course access", "onboarding materials"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "analytics", "disengagement", "intervention",
                        "quiz", "checklist", "advisor", "workshop"
                    },
                    "phrases": [
                        "intelligent agents", "early alerts",
                        "class progress", "engagement dashboards",
                        "active learning", "digital engagement",
                        "engagement reports"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "predictive", "retention", "progression",
                        "holistic", "inclusive", "scalable"
                    },
                    "phrases": [
                        "predictive analytics", "student success",
                        "data hub", "cross-functional initiatives",
                        "institutional goals", "real-time engagement",
                        "students at risk"
                    ]
                }
            }
        },
        "04_data_analytics": {
            "name": "Data & Learning Analytics",
            "short": "Data Analytics",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "report", "export", "gradebook", "basic",
                        "manual", "spreadsheet"
                    },
                    "phrases": [
                        "class progress", "basic reports",
                        "grade exports", "manual tracking",
                        "default reports"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "analytics", "dashboard", "visualization",
                        "benchmark", "kpi", "metric", "trend"
                    },
                    "phrases": [
                        "data-informed", "engagement dashboards",
                        "learning analytics", "custom reports",
                        "data hub", "decision making"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "predictive", "modeling", "warehouse",
                        "longitudinal", "algorithm"
                    },
                    "phrases": [
                        "predictive analytics", "power bi", "tableau",
                        "data warehouse", "institutional research",
                        "cross-functional initiatives",
                        "data governance"
                    ]
                }
            }
        },
        "05_assessment": {
            "name": "Assessment & Evaluation",
            "short": "Assessment",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "quiz", "assignment", "gradebook", "summative",
                        "multiple-choice", "grading", "manual"
                    },
                    "phrases": [
                        "auto-graded quizzes", "default settings",
                        "assignment uploads", "basic assignments",
                        "getting started", "step-by-step tutorials"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "rubric", "formative", "peer", "reflective",
                        "randomization", "integrity", "feedback"
                    },
                    "phrases": [
                        "peer review", "question banks",
                        "learning outcomes", "peer assessment",
                        "assessment design", "rubric libraries",
                        "knowledge checks", "video assignments"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "adaptive", "simulation", "personalized",
                        "branching", "turnitin", "h5p", "innovation"
                    },
                    "phrases": [
                        "item analysis", "real-time feedback",
                        "adaptive learning", "innovation labs",
                        "pilot projects", "assessment innovation",
                        "delegated marking", "co-marking"
                    ]
                }
            }
        },
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
        "06_instructor_efficiency": {
            "name": "Instructor Efficiency",
            "short": "Instructor Efficiency",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "onboarding", "reactive", "tutorial", "guide",
                        "slides", "compliance", "upload", "training",
                        "support", "help", "manual", "orientation",
                        "faculty", "instructor"
                    },
                    "phrases": [
                        "step-by-step guides", "video tutorials",
                        "getting started", "drop-in help",
                        "technical support", "support dependency",
                        "faculty training", "instructor support",
                        "basic training", "how-to guides"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "proactive", "contextual", "champion", "proficiency",
                        "troubleshoot", "multimedia", "padlet", "jamboard",
                        "captioned", "coaching", "workshop", "consultation",
                        "adoption", "readiness", "enablement"
                    },
                    "phrases": [
                        "best practices", "faculty champions",
                        "pedagogical integration", "analytics dashboards",
                        "peer learning", "faculty showcases",
                        "user groups", "one-on-one coaching",
                        "faculty development", "contextual support",
                        "training sessions", "teaching support"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "collaborative", "strategic", "innovation",
                        "mentorship", "gamified", "flipped", "adaptive",
                        "personalized", "agency", "transformation",
                        "scalable", "empowerment"
                    },
                    "phrases": [
                        "change agents", "custom tool integrations",
                        "data-driven", "beta testing",
                        "ai tools", "edtech innovation",
                        "data exports", "continuous improvement",
                        "faculty learning communities",
                        "instructional design partnership",
                        "teaching innovation"
                    ]
                }
            }
        },
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<
        "07_change_management": {
            "name": "Change Management",
            "short": "Change Mgmt",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "update", "outage", "rollout", "disruption",
                        "communicate", "documentation"
                    },
                    "phrases": [
                        "system updates", "new feature",
                        "just-in-time support", "minimize disruption",
                        "pros and cons"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "adoption", "sandbox", "webinar", "feedback",
                        "champion", "training"
                    },
                    "phrases": [
                        "training sessions", "sandbox environments",
                        "rollout strategies", "instructor feedback",
                        "best practices", "faculty champions"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "transformation", "strategy", "innovation",
                        "pilot", "culture"
                    },
                    "phrases": [
                        "change management plans", "digital transformation",
                        "continuous improvement", "communities of practice",
                        "cross-functional collaboration",
                        "pilot programs"
                    ]
                }
            }
        },
        "08_knowledge_management": {
            "name": "Knowledge & Resource Management",
            "short": "Knowledge Mgmt",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "documentation", "faq", "guide", "reactive",
                        "manual", "organize", "upload"
                    },
                    "phrases": [
                        "knowledge base", "help site",
                        "how-to guides", "support channels",
                        "course shells", "case-by-case"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "template", "library", "curate", "repository",
                        "lor", "reuse", "redundancy", "analytics"
                    },
                    "phrases": [
                        "shared content", "content libraries",
                        "course copy", "instructional designers",
                        "institutional standards", "content usage"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "governance", "lifecycle", "versioning",
                        "archiving", "metadata", "tagging",
                        "discoverability", "stewardship"
                    },
                    "phrases": [
                        "communities of practice", "knowledge strategy",
                        "content lifecycle", "metadata standards",
                        "cross-departmental collaboration",
                        "oer repositories"
                    ]
                }
            }
        },
        "09_accessibility": {
            "name": "Accessibility & Compliance",
            "short": "Accessibility",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "accessibility", "ada", "wcag", "captioning",
                        "keyboard", "alt", "readspeaker", "ally"
                    },
                    "phrases": [
                        "screen reader", "alternative text",
                        "accessibility standards", "colour contrast",
                        "accessible content", "legal requirements"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "audit", "remediation", "inclusive", "assistive",
                        "consultation", "template", "checklist"
                    },
                    "phrases": [
                        "inclusive course design", "assistive technologies",
                        "course accessibility", "accessible course materials",
                        "content remediation", "captioning legacy"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "champion", "mentor", "procurement", "innovation",
                        "immersive"
                    },
                    "phrases": [
                        "accessibility champions", "accessibility compliance",
                        "procurement policies", "ai-driven captioning",
                        "inclusive pedagogy", "accessibility audits"
                    ]
                }
            }
        },
        "10_user_support": {
            "name": "User Support & Training",
            "short": "User Support",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "login", "navigation", "tutorial", "faq",
                        "guide", "gradebook", "assignment"
                    },
                    "phrases": [
                        "quick-start guides", "video tutorials",
                        "login support", "basic training",
                        "core tools"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "tiered", "workshop", "rubric", "multimedia",
                        "contextualized"
                    },
                    "phrases": [
                        "tiered training", "course design",
                        "assessment strategies", "student engagement",
                        "release conditions", "contextualized support"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "coaching", "consultation", "champion",
                        "personalize", "competency"
                    },
                    "phrases": [
                        "one-on-one coaching", "instructional design consultations",
                        "faculty learning communities",
                        "intelligent agents", "competency frameworks",
                        "learning analytics"
                    ]
                }
            }
        },
        "11_innovation": {
            "name": "Innovation & Emerging Technologies",
            "short": "Innovation",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "sandbox", "demo", "showcase", "curiosity",
                        "explore", "experiment", "gamification"
                    },
                    "phrases": [
                        "sandbox environments", "demo sessions",
                        "tech showcases", "immersive media",
                        "ai-powered tools"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "pilot", "adaptive", "simulation", "framework",
                        "evaluate", "effectiveness"
                    },
                    "phrases": [
                        "hands-on training", "pilot programs",
                        "adaptive learning", "ai feedback",
                        "virtual simulations", "technology adoption"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "research", "innovation", "equity", "transformation",
                        "scaled"
                    },
                    "phrases": [
                        "innovation labs", "communities of practice",
                        "digital pedagogy", "institutional goals",
                        "change agents", "faculty-led research"
                    ]
                }
            }
        },
        "12_collaboration": {
            "name": "Collaboration & Communication",
            "short": "Collaboration",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "support", "troubleshooting", "reactive",
                        "coordinate", "respond"
                    },
                    "phrases": [
                        "support requests", "system updates",
                        "service platform", "issue resolution",
                        "one-way communication", "tool focused"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "consultative", "engage", "curriculum",
                        "faculty", "feedback", "align"
                    },
                    "phrases": [
                        "pedagogical partner", "two-way communication",
                        "curriculum planning", "faculty development",
                        "cross-functional meetings",
                        "academic priorities"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "strategic", "governance", "transformation",
                        "vision", "proactive"
                    },
                    "phrases": [
                        "institutional strategy", "governance committees",
                        "strategic initiatives", "digital transformation",
                        "shared vision", "data-informed",
                        "emerging technologies"
                    ]
                }
            }
        }
    }
}
        }

    def get_model_names(self) -> List[str]:
        return list(self.models.keys())

    def get_model_desc(self, name: str) -> str:
        return self.models.get(name, {}).get("desc", "")

    def assess(self, counts: Counter, bigrams: Counter, model_name: str) -> Dict:
        """Routes to flat or domain-based assessment based on model type."""
        if model_name not in self.models:
            return None
        model = self.models[model_name]
        if model.get("type") == "domain_based":
            return self._assess_domain_based(counts, bigrams, model_name)
        else:
            return self._assess_flat(counts, bigrams, model_name)
    def _assess_flat(self, counts: Counter, bigrams: Counter, model_name: str) -> Dict:
        """Original 5-level flat assessment. Unchanged logic."""
        if model_name not in self.models:
            return None
        levels = self.models[model_name]["levels"]
        scores = {1: 0.0, 2: 0.0, 3: 0.0, 4: 0.0, 5: 0.0}
        total_hits = 0.0
        for lvl, data in levels.items():
            terms = data.get("terms", set())
            for term in terms:
                qty = counts.get(term, 0)
                if qty > 0:
                    scores[lvl] += qty
                    total_hits += qty
        phrase_counts: Dict[str, int] = {}
        for (w1, w2), freq in bigrams.items():
            phrase = f"{w1} {w2}"
            phrase_counts[phrase] = phrase_counts.get(phrase, 0) + freq
        PHRASE_WEIGHT = 1.5
        for lvl, data in levels.items():
            phrases = data.get("phrases", []) or []
            for phrase in phrases:
                freq = phrase_counts.get(phrase, 0)
                if freq > 0:
                    weighted = freq * PHRASE_WEIGHT
                    scores[lvl] += weighted
                    total_hits += weighted
        if total_hits == 0:
            return None
        distribution = {k: v / total_hits for k, v in scores.items()}
        weighted_sum = sum(lvl * pct for lvl, pct in distribution.items())
        dominant_level = max(distribution, key=distribution.get)
        return {
            "type": "flat",
            "overall_score": round(weighted_sum, 2),
            "distribution": distribution,
            "dominant_stage": levels[dominant_level],
            "total_signals_found": int(total_hits),
            "levels_ref": levels
        }
    def _assess_domain_based(self, counts: Counter, bigrams: Counter, model_name: str) -> Dict:
        """
        12-domain, 3-tier assessment. Returns per-domain scores + composite.
        Score range: 1.0 (all Foundational) to 3.0 (all Leading Edge).
        """
        if model_name not in self.models:
            return None
        domains = self.models[model_name]["domains"]
        PHRASE_WEIGHT = 1.5
        # Build bigram lookup once
        phrase_counts: Dict[str, int] = {}
        for (w1, w2), freq in bigrams.items():
            phrase = f"{w1} {w2}"
            phrase_counts[phrase] = phrase_counts.get(phrase, 0) + freq
        domain_results = {}
        total_signals_all = 0
        for domain_key, domain_data in domains.items():
            tiers = domain_data["tiers"]
            tier_scores = {1: 0.0, 2: 0.0, 3: 0.0}
            domain_hits = 0.0
            driver_terms = {}
            for tier_num, tier_data in tiers.items():
                tier_term_hits = []
                # Unigram hits
                for term in tier_data.get("terms", set()):
                    qty = counts.get(term, 0)
                    if qty > 0:
                        tier_scores[tier_num] += qty
                        domain_hits += qty
                        tier_term_hits.append((term, qty))
                # Phrase hits
                for phrase in tier_data.get("phrases", []):
                    freq = phrase_counts.get(phrase, 0)
                    if freq > 0:
                        weighted = freq * PHRASE_WEIGHT
                        tier_scores[tier_num] += weighted
                        domain_hits += weighted
                        tier_term_hits.append((phrase, freq))
                # Sort by count descending, keep top 5
                driver_terms[tier_num] = sorted(
                    tier_term_hits, key=lambda x: x[1], reverse=True
                )[:5]
            if domain_hits == 0:
                domain_results[domain_key] = {
                    "name": domain_data["name"],
                    "short": domain_data["short"],
                    "score": 0.0,
                    "tier_label": "No Data",
                    "distribution": {1: 0.0, 2: 0.0, 3: 0.0},
                    "signals": 0,
                    "drivers": driver_terms
                }
                continue
            total_signals_all += domain_hits
            # Normalize distribution
            distribution = {k: v / domain_hits for k, v in tier_scores.items()}
            # Weighted score (1.0 - 3.0)
            weighted_score = sum(tier * pct for tier, pct in distribution.items())
            # Determine tier label
            if weighted_score >= 2.5:
                tier_label = "Leading Edge"
            elif weighted_score >= 1.5:
                tier_label = "Advanced"
            else:
                tier_label = "Foundational"
            domain_results[domain_key] = {
                "name": domain_data["name"],
                "short": domain_data["short"],
                "score": round(weighted_score, 2),
                "tier_label": tier_label,
                "distribution": distribution,
                "signals": int(domain_hits),
                "drivers": driver_terms
            }
        # Composite score (average of domains that have data)
        scored_domains = [
            d for d in domain_results.values() if d["signals"] > 0
        ]
        if not scored_domains:
            return None
        composite_score = sum(d["score"] for d in scored_domains) / len(scored_domains)
        return {
            "type": "domain_based",
            "overall_score": round(composite_score, 2),
            "max_score": 3.0,
            "domains_assessed": len(scored_domains),
            "domains_total": 12,
            "total_signals_found": int(total_signals_all),
            "domain_results": domain_results,
            "domains_ref": domains
        }

    def render_radar_chart(self, result: Dict):
        """Generates a Radar/Spider chart."""
        if not result:
            return None

        levels_ref = result['levels_ref']
        categories = [levels_ref[i]['name'] for i in range(1, 6)]
        values = [result['distribution'][i] for i in range(1, 6)]

        values += values[:1]
        angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
        angles += angles[:1]

        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
        ax.fill(angles, values, color='#1f77b4', alpha=0.25)
        ax.plot(angles, values, color='#1f77b4', linewidth=2)
        ax.set_theta_offset(np.pi / 2)
        ax.set_theta_direction(-1)
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(categories, fontsize=9)
        ax.set_yticklabels([])
        ax.spines["polar"].set_visible(False)
        ax.grid(color='#444444', linestyle='--', alpha=0.5)
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)

        return fig
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
    def render_domain_radar_chart(
        self,
        result: Dict,
        font_family: str = "sans-serif",
        label_size: int = 9,
        label_color: str = "#222222",
        wrap_width: int = 14,
        show_tier_labels: bool = True,
        tier_label_angle: int = 35,
    ):
        """12-spoke spider chart for domain-based assessment."""
        if not result or result.get("type") != "domain_based":
            return None

        import textwrap

        domain_results = result["domain_results"]
        labels = []
        values = []
        colors = []

        for key in sorted(domain_results.keys()):
            dr = domain_results[key]
            wrapped_label = "\n".join(
                textwrap.wrap(str(dr["short"]), width=max(6, wrap_width))
            )
            labels.append(wrapped_label)
            values.append(dr["score"])

            if dr["score"] >= 2.5:
                colors.append("#2ca02c")
            elif dr["score"] >= 1.5:
                colors.append("#ff7f0e")
            else:
                colors.append("#d62728")

        n = len(labels)
        if n == 0:
            return None

        values_closed = values + values[:1]
        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        angles_closed = angles + angles[:1]

        fig, ax = plt.subplots(figsize=(8.2, 8.2), subplot_kw=dict(polar=True))
        ax.fill(angles_closed, values_closed, color="#1f77b4", alpha=0.15)
        ax.plot(angles_closed, values_closed, color="#1f77b4", linewidth=2)

        for angle, val, color in zip(angles, values, colors):
            ax.scatter(angle, val, color=color, s=80, zorder=5)

        ax.set_theta_offset(np.pi / 2)
        ax.set_theta_direction(-1)
        ax.set_xticks(angles)
        ax.set_xticklabels(
            labels,
            fontsize=label_size,
            fontfamily=font_family,
            color=label_color,
        )

        # Move domain labels outward to reduce collisions with ring labels.
        ax.tick_params(axis="x", pad=18)

        ax.set_ylim(0, 3.0)
        ax.set_yticks([1.0, 2.0, 3.0])

        if show_tier_labels:
            ax.set_yticklabels(
                ["Foundational", "Advanced", "Leading\nEdge"],
                fontsize=max(7, label_size - 2),
                fontfamily=font_family,
                color=label_color,
            )
        else:
            ax.set_yticklabels(
                ["1.0", "2.0", "3.0"],
                fontsize=max(7, label_size - 2),
                fontfamily=font_family,
                color=label_color,
            )

        # Move maturity ring labels away from the crowded top-right label area.
        ax.set_rlabel_position(tier_label_angle)

        ax.spines["polar"].set_visible(False)
        ax.grid(color="#444444", linestyle="--", alpha=0.5)
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        fig.tight_layout(pad=2.5)

        return fig
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<
        
        
    def render_domain_breakdown_chart(self, result: Dict):
        """Horizontal stacked bar chart showing tier distribution per domain."""
        if not result or result.get("type") != "domain_based":
            return None
        domain_results = result["domain_results"]
        sorted_keys = sorted(domain_results.keys())
        labels = []
        found_vals = []
        adv_vals = []
        lead_vals = []
        for key in sorted_keys:
            dr = domain_results[key]
            labels.append(dr["short"])
            dist = dr["distribution"]
            found_vals.append(dist.get(1, 0) * 100)
            adv_vals.append(dist.get(2, 0) * 100)
            lead_vals.append(dist.get(3, 0) * 100)
        y_pos = np.arange(len(labels))
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.barh(y_pos, found_vals, color='#d62728', label='Foundational', height=0.6)
        ax.barh(y_pos, adv_vals, left=found_vals, color='#ff7f0e', label='Advanced', height=0.6)
        lefts = [f + a for f, a in zip(found_vals, adv_vals)]
        ax.barh(y_pos, lead_vals, left=lefts, color='#2ca02c', label='Leading Edge', height=0.6)
        ax.set_yticks(y_pos)
        ax.set_yticklabels(labels, fontsize=9)
        ax.set_xlabel("Signal Distribution (%)")
        ax.set_xlim(0, 100)
        ax.legend(loc='lower right', fontsize=8)
        ax.invert_yaxis()
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        ax.grid(axis='x', alpha=0.3)
        plt.tight_layout()
        return fig
    def render_trend_chart(self, snapshots: list):
        """Line chart showing composite score over time across snapshots."""
        if not snapshots or len(snapshots) < 2:
            return None
        dates = [s["assessment_date"] for s in snapshots]
        scores = [s["composite_score"] for s in snapshots]
        fig, ax = plt.subplots(figsize=(10, 4))
        ax.plot(dates, scores, marker='o', linewidth=2, color='#1f77b4', markersize=8)
        # Color-code background bands
        ax.axhspan(1.0, 1.5, alpha=0.08, color='#d62728', label='Foundational')
        ax.axhspan(1.5, 2.5, alpha=0.08, color='#ff7f0e', label='Advanced')
        ax.axhspan(2.5, 3.0, alpha=0.08, color='#2ca02c', label='Leading Edge')
        ax.set_ylim(0.8, 3.2)
        ax.set_ylabel("Composite Maturity Score")
        ax.set_xlabel("Assessment Date")
        ax.legend(loc='lower right', fontsize=8)
        plt.xticks(rotation=45, ha='right')
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        ax.grid(axis='y', alpha=0.3)
        plt.tight_layout()
        return fig
    def render_domain_trend_chart(self, snapshots: list):
        """Multi-line chart showing per-domain scores over time."""
        if not snapshots or len(snapshots) < 2:
            return None
        dates = [s["assessment_date"] for s in snapshots]
        # Collect all domain keys across all snapshots
        all_domains = set()
        for s in snapshots:
            all_domains.update(s.get("domain_results", {}).keys())
        fig, ax = plt.subplots(figsize=(12, 6))
        domain_colors = [
            "#e6194b", "#3cb44b", "#ffe119", "#4363d8", "#f58231", "#911eb4",
            "#42d4f4", "#f032e6", "#bfef45", "#fabed4", "#469990", "#dcbeff"
        ]
        for i, dk in enumerate(sorted(all_domains)):
            domain_scores = []
            for s in snapshots:
                dr = s.get("domain_results", {}).get(dk, {})
                domain_scores.append(dr.get("score", 0.0))
            short_name = snapshots[0].get("domain_results", {}).get(
                dk, {}
            ).get("short", dk)
            color = domain_colors[i % len(domain_colors)]
            ax.plot(dates, domain_scores, marker='o', linewidth=1.5,
                    label=short_name, color=color, markersize=5)
        ax.set_ylim(0.8, 3.2)
        ax.set_ylabel("Domain Score")
        ax.set_xlabel("Assessment Date")
        ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', fontsize=7)
        plt.xticks(rotation=45, ha='right')
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        ax.grid(axis='y', alpha=0.3)
        plt.tight_layout()
        return fig


# session state init
if 'sketch' not in st.session_state: st.session_state['sketch'] = StreamScanner()
if 'total_cost' not in st.session_state: st.session_state['total_cost'] = 0.0
if 'total_tokens' not in st.session_state: st.session_state['total_tokens'] = 0
if 'authenticated' not in st.session_state: st.session_state['authenticated'] = False
if 'auth_error' not in st.session_state: st.session_state['auth_error'] = False
if 'ai_response' not in st.session_state: st.session_state['ai_response'] = ""
if 'last_sketch_hash' not in st.session_state: st.session_state['last_sketch_hash'] = None

def reset_sketch():
    st.session_state['sketch'] = StreamScanner()
    st.session_state['ai_response'] = ""
    st.session_state['last_sketch_hash'] = None
    gc.collect()

def perform_login():
    try:
        correct_password = get_auth_password()
        if secrets.compare_digest(st.session_state.password_input, correct_password):
            st.session_state['authenticated'] = True
            st.session_state['auth_error'] = False
            st.session_state['password_input'] = ""
        else:
            st.session_state['auth_error'] = True
    except Exception:
        st.session_state['auth_error'] = True

def logout():
    st.session_state['authenticated'] = False
    st.session_state['ai_response'] = ""


# 🛠️ helpers, setup
# ============================================

@st.cache_resource(show_spinner="Init NLTK...")
def setup_nlp_resources():
    if nltk is None: return None, None
    try: 
        nltk.data.find('sentiment/vader_lexicon.zip')
        nltk.data.find('corpora/wordnet.zip')
        nltk.data.find('corpora/omw-1.4.zip')
    except LookupError: 
        try:
            nltk.download('vader_lexicon')
            nltk.download('wordnet')
            nltk.download('omw-1.4')
        except:
            pass
    
    sia = SentimentIntensityAnalyzer()
    lemmatizer = WordNetLemmatizer()
    return sia, lemmatizer

@st.cache_data(show_spinner=False)
def list_system_fonts() -> Dict[str, str]:
    mapping = {}
    for fe in font_manager.fontManager.ttflist:
        if fe.name not in mapping: mapping[fe.name] = fe.fname
    return dict(sorted(mapping.items(), key=lambda x: x[0].lower()))

def build_punct_translation(keep_hyphens: bool, keep_apostrophes: bool) -> dict:
    # 1. Standard ASCII punctuation
    punct = string.punctuation
    
    # 2. Add Unicode "Smart" quotes & dashes
    punct += "“”‘’–—" 

    if keep_hyphens: 
        for char in "-–—": punct = punct.replace(char, "")
    if keep_apostrophes: 
        for char in "'’": punct = punct.replace(char, "")
    
    return str.maketrans("", "", punct)

def parse_user_stopwords(raw: str) -> Tuple[List[str], List[str]]:
    raw = raw.replace("\n", ",").replace(".", ",")
    phrases, singles = [], []
    for item in [x.strip() for x in raw.split(",") if x.strip()]:
        if " " in item: phrases.append(item.lower())
        else: singles.append(item.lower())
    return phrases, singles

#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
def default_prepositions() -> set:
    return {
        'about', 'above', 'across', 'after', 'against', 'along', 'among',
        'around', 'at', 'before', 'behind', 'below', 'beneath', 'beside',
        'between', 'beyond', 'but', 'by', 'concerning', 'despite', 'down',
        'during', 'except', 'for', 'from', 'in', 'inside', 'into', 'like',
        'near', 'of', 'off', 'on', 'onto', 'out', 'outside', 'over', 'past',
        'regarding', 'since', 'through', 'throughout', 'to', 'toward', 'under',
        'underneath', 'until', 'up', 'upon', 'with', 'within', 'without',
        'something', 'someone', 'somebody', 'anything', 'anyone', 'anybody',
        'everything', 'everyone', 'everybody', 'nothing', 'none', 'maybe',
        'perhaps', 'really', 'basically', 'generally', 'kind', 'sort', 'stuff',
        'thing', 'things'
    }
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<

#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
def build_phrase_pattern(phrases: List[str]) -> Optional[re.Pattern]:
    if not phrases:
        return None
    escaped = [re.escape(p) for p in phrases if p]
    if not escaped:
        return None
    return re.compile(rf"\b(?:{'|'.join(escaped)})\b", flags=re.IGNORECASE)


def collect_maturity_vocabulary() -> Set[str]:
    """
    Collects maturity-model vocabulary so maturity-critical words are not
    accidentally removed by generic stopword / filler-word cleanup.

    This intentionally includes:
    - explicit one-word terms
    - individual words from maturity phrases

    Example: "student success" protects both "student" and "success".
    """
    protected_terms: Set[str] = set()

    try:
        assessor = MaturityAssessor()
        for model in assessor.models.values():
            if model.get("type") == "domain_based":
                domain_iter = model.get("domains", {}).values()
                for domain_data in domain_iter:
                    for tier_data in domain_data.get("tiers", {}).values():
                        for term in tier_data.get("terms", set()):
                            protected_terms.add(str(term).lower())
                        for phrase in tier_data.get("phrases", []):
                            protected_terms.update(
                                part.lower()
                                for part in str(phrase).split()
                                if part.strip()
                            )
            else:
                for level_data in model.get("levels", {}).values():
                    for term in level_data.get("terms", set()):
                        protected_terms.add(str(term).lower())
                    for phrase in level_data.get("phrases", []):
                        protected_terms.update(
                            part.lower()
                            for part in str(phrase).split()
                            if part.strip()
                        )
    except Exception:
        # Safety-first fallback: never break app startup because help/protection
        # vocabulary could not be collected.
        return set()

    return protected_terms


def protect_maturity_vocabulary(stopwords: Set[str]) -> Set[str]:
    """
    Removes maturity-critical terms from the active stopword set.

    This preserves maturity scoring while still allowing generic filler words
    and user-entered junk terms to clean the broader analysis.
    """
    maturity_terms = collect_maturity_vocabulary()
    if not maturity_terms:
        return stopwords
    return set(stopwords) - maturity_terms
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<

def estimate_row_count_from_bytes(file_bytes: bytes) -> int:
    if not file_bytes: return 0
    return file_bytes.count(b'\n') + 1

def make_unique_header(raw_names: List[Optional[str]]) -> List[str]:
    seen: Dict[str, int] = {}
    result: List[str] = []
    for i, nm in enumerate(raw_names):
        name = (str(nm).strip() if nm is not None else "")
        if not name: name = f"col_{i}"
        if name in seen:
            seen[name] += 1
            unique = f"{name}__{seen[name]}"
        else:
            seen[name] = 1
            unique = name
        result.append(unique)
    return result

def extract_entities_regex(text: str, stopwords: Set[str]) -> List[str]:
    # lightweight NER without heavy models
    candidates = NER_CAPS_RE.findall(text)
    valid = []
    for c in candidates:
        # filter out if it's just a common stopword capitalized at start of sentence
        if c.lower() in stopwords: continue
        if len(c) < 3: continue
        valid.append(c)
    return valid

# --- virtual files and web/url
class VirtualFile:
    def __init__(self, name: str, text_content: str):
        self.name = name
        self._bytes = text_content.encode('utf-8')
    
    def getvalue(self) -> bytes:
        return self._bytes
    
    def getbuffer(self) -> memoryview:
        return memoryview(self._bytes)

def fetch_url_content(url: str) -> Optional[str]:
    if not requests or not BeautifulSoup: return None
    if not validate_url(url): return None
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        for script in soup(["script", "style", "nav", "footer"]):
            script.decompose()
        return soup.get_text(separator=' ', strip=True)
    except Exception: return None


# 📄 file readers (tuple yielding)
# ==========================================

# all readers yield (text_content, date_str, category_str)

def read_rows_raw_lines(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[Tuple[str, None, None]]:
    def _iter(enc):
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline=None) as wrapper:
            for line in wrapper: yield (line.rstrip("\r\n"), None, None)
    try:
        if encoding_choice == "latin-1": yield from _iter("latin-1")
        else: yield from _iter("utf-8")
    except UnicodeDecodeError:
        yield ("", None, None)

def read_rows_vtt(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[Tuple[str, None, None]]:
    # robust VTT reader that yields tuples
    def _iter_lines(enc):
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline=None) as wrapper:
            for line in wrapper: yield line.rstrip("\r\n")
    
    iterator = _iter_lines("utf-8") if encoding_choice != "latin-1" else _iter_lines("latin-1")
    
    for line in iterator:
        line = line.strip()
        if not line or line == "WEBVTT" or "-->" in line or line.isdigit(): continue
        if ":" in line:
            parts = line.split(":", 1)
            if len(parts) > 1 and len(parts[0]) < MAX_SPEAKER_NAME_LENGTH and " " in parts[0]:
                yield (parts[1].strip(), None, None)
                continue
        yield (line, None, None)

def read_rows_pdf(file_bytes: bytes) -> Iterable[Tuple[str, None, None]]:
    if pypdf is None: 
        st.error("pypdf missing")
        return
    bio = io.BytesIO(file_bytes)
    try:
        reader = pypdf.PdfReader(bio)
        for page in reader.pages:
            text = page.extract_text()
            if text: yield (text, None, None)
    except Exception: yield ("", None, None)

def read_rows_pptx(file_bytes: bytes) -> Iterable[Tuple[str, None, None]]:
    if pptx is None: return
    bio = io.BytesIO(file_bytes)
    try:
        prs = pptx.Presentation(bio)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.text: yield (shape.text, None, None)
    except Exception:
        yield ("", None, None)

def read_rows_json(file_bytes: bytes, selected_key: str = None) -> Iterable[Tuple[str, None, None]]:
    bio = io.BytesIO(file_bytes)
    try:
        wrapper = io.TextIOWrapper(bio, encoding="utf-8", errors="replace")
        for line in wrapper:
            if not line.strip(): continue
            try:
                obj = json.loads(line)
                txt = ""
                if selected_key and isinstance(obj, dict): txt = str(obj.get(selected_key, ""))
                elif isinstance(obj, str): txt = obj
                else: txt = str(obj)
                yield (txt, None, None)
            except json.JSONDecodeError:
                pass
    except Exception:
        pass

def read_rows_csv_structured(
    file_bytes: bytes, 
    encoding_choice: str, 
    delimiter: str, 
    has_header: bool, 
    text_cols: List[str],
    date_col: Optional[str],
    cat_col: Optional[str],
    join_with: str
) -> Iterable[Tuple[str, Optional[str], Optional[str]]]:
    
    enc = "latin-1" if encoding_choice == "latin-1" else "utf-8"
    bio = io.BytesIO(file_bytes)
    
    with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline="") as wrapper:
        rdr = csv.reader(wrapper, delimiter=delimiter)
        first = next(rdr, None)
        if first is None: return

        if has_header:
            header = make_unique_header(list(first))
            name_to_idx = {n: i for i, n in enumerate(header)}
            
            text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
            date_idx = name_to_idx.get(date_col) if date_col else None
            cat_idx = name_to_idx.get(cat_col) if cat_col else None
        else:
            # if no header, user likely selected "col_0", "col_1" etc.
            name_to_idx = {f"col_{i}": i for i in range(len(first))}
            text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
            date_idx = name_to_idx.get(date_col) if date_col else None
            cat_idx = name_to_idx.get(cat_col) if cat_col else None
            
            # yield 1st row data
            txt_parts = [first[i] if i < len(first) else "" for i in text_idxs]
            d_val = first[date_idx] if (date_idx is not None and date_idx < len(first)) else None
            c_val = first[cat_idx] if (cat_idx is not None and cat_idx < len(first)) else None
            yield (join_with.join(txt_parts), d_val, c_val)

        for row in rdr:
            txt_parts = [row[i] if i < len(row) else "" for i in text_idxs]
            d_val = row[date_idx] if (date_idx is not None and date_idx < len(row)) else None
            c_val = row[cat_idx] if (cat_idx is not None and cat_idx < len(row)) else None
            yield (join_with.join(txt_parts), d_val, c_val)

def iter_excel_structured(
    file_bytes: bytes, 
    sheet_name: str, 
    has_header: bool, 
    text_cols: List[str],
    date_col: Optional[str],
    cat_col: Optional[str],
    join_with: str
) -> Iterable[Tuple[str, Optional[str], Optional[str]]]:
    if openpyxl is None: return
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    ws = wb[sheet_name]
    rows_iter = ws.iter_rows(values_only=True)
    
    first = next(rows_iter, None)
    if first is None: 
        wb.close()
        return

    # header logic
    if has_header:
        header = make_unique_header(list(first))
        name_to_idx = {n: i for i, n in enumerate(header)}
        text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
        date_idx = name_to_idx.get(date_col) if date_col else None
        cat_idx = name_to_idx.get(cat_col) if cat_col else None
    else:
        name_to_idx = {f"col_{i}": i for i in range(len(first))}
        text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
        date_idx = name_to_idx.get(date_col) if date_col else None
        cat_idx = name_to_idx.get(cat_col) if cat_col else None
        
        # Yield first row
        txt_parts = [str(first[i]) if (first[i] is not None and i < len(first)) else "" for i in text_idxs]
        d_val = first[date_idx] if (date_idx is not None and date_idx < len(first)) else None
        c_val = first[cat_idx] if (cat_idx is not None and cat_idx < len(first)) else None
        yield (join_with.join(txt_parts), str(d_val) if d_val else None, str(c_val) if c_val else None)

    for row in rows_iter:
        txt_parts = [str(row[i]) if (row[i] is not None and i < len(row)) else "" for i in text_idxs]
        d_val = row[date_idx] if (date_idx is not None and date_idx < len(row)) else None
        c_val = row[cat_idx] if (cat_idx is not None and cat_idx < len(row)) else None
        yield (join_with.join(txt_parts), str(d_val) if d_val else None, str(c_val) if c_val else None)
    
    wb.close()

def detect_csv_headers(file_bytes: bytes, delimiter: str = ",") -> List[str]:
    try:
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding="utf-8", errors="replace", newline="") as wrapper:
            rdr = csv.reader(wrapper, delimiter=delimiter)
            row = next(rdr, None)
            return make_unique_header(row) if row else []
    except: return []

def detect_csv_num_cols(file_bytes: bytes, delimiter: str = ",") -> int:
    try:
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding="utf-8", errors="replace", newline="") as wrapper:
            rdr = csv.reader(wrapper, delimiter=delimiter)
            row = next(rdr, None)
            return len(row) if row else 0
    except: return 0

def get_excel_sheetnames(file_bytes: bytes) -> List[str]:
    if openpyxl is None: return []
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    sheets = list(wb.sheetnames)
    wb.close()
    return sheets

def get_excel_preview(file_bytes: bytes, sheet_name: str, has_header: bool, rows: int = 5) -> pd.DataFrame:
    if openpyxl is None: return pd.DataFrame()
    bio = io.BytesIO(file_bytes)
    try:
        df = pd.read_excel(bio, sheet_name=sheet_name, header=0 if has_header else None, nrows=rows, engine='openpyxl')
        if not has_header: df.columns = [f"col_{i}" for i in range(len(df.columns))]
        return df
    except:
        return pd.DataFrame()

def excel_estimate_rows(file_bytes: bytes, sheet_name: str, has_header: bool) -> int:
    if openpyxl is None: return 0
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    ws = wb[sheet_name]
    total = ws.max_row or 0
    wb.close()
    if has_header and total > 0: total -= 1
    return max(total, 0)


# ⚙️ processing logic
# ==========================================
def clean_date_str(raw: Any) -> Optional[str]:
    """Robust extraction using dateutil if available, falling back to regex."""
    if not raw: return None
    s = str(raw).strip()
    
    # tries smart parsing (covers "Jan 5, 2024", "2024/01/01", etc)
    if date_parser:
        try:
            # fuzzy=True allows extracting dates buried in strings like "Date: Jan 1"
            dt = date_parser.parse(s, fuzzy=True)
            return dt.strftime("%Y-%m-%d")
        except (ValueError, OverflowError):
            pass # falls through to regex if dateutil fails

    # the fallback Regex (simple ISO & US formats)
    match = re.search(r'\d{4}-\d{2}-\d{2}', s)
    if match: return match.group(0)
    match_us = re.search(r'(\d{1,2})/(\d{1,2})/(\d{4})', s)
    if match_us:
        m, d, y = match_us.groups()
        return f"{y}-{int(m):02d}-{int(d):02d}"
    
    return None

def apply_text_cleaning(text: str, config: CleaningConfig) -> str:
    if not isinstance(text, str): return ""
    
    # standard cleaning FIRST (to convert tags/artifacts to spaces)
    # ensures artifacts like <br> become spaces so the de-hyphenator can see them
    if config.remove_chat: text = CHAT_ARTIFACT_RE.sub(" ", text)
    if config.remove_html: text = HTML_TAG_RE.sub(" ", text)
    if config.unescape:
        try: text = html.unescape(text)
        except: pass

    # the de-hyphenation repair
    # merges words split by hyphen+whitespace (e.g., "equiv-\nalent" -> "equivalent")
    text = re.sub(r'(\w)-\s+(\w)', r'\1\2', text)
    
    # sticky Acronyms Fix (e.g., "WHOrecommended" -> "WHO recommended")
    # logic: 2+ CAPS followed by 2+ lowercase letters
    text = re.sub(r'([A-Z]{2,})([a-z]{2,})', r'\1 \2', text)

    # final cleanups
    if config.remove_urls: text = URL_EMAIL_RE.sub(" ", text)
    
    # normalizing
    text = text.lower()
    
    if config.phrase_pattern: text = config.phrase_pattern.sub(" ", text)
    return text.strip()

def process_chunk_iter(
    rows_iter: Iterable[Tuple[str, Optional[str], Optional[str]]],
    clean_conf: CleaningConfig,
    proc_conf: ProcessingConfig,
    scanner: StreamScanner,
    lemmatizer: Optional[WordNetLemmatizer],
    progress_cb: Optional[Callable[[int], None]] = None
):
    _min_len = proc_conf.min_word_len
    _drop_int = proc_conf.drop_integers
    _trans = proc_conf.translate_map
    _stopwords = proc_conf.stopwords
    _lemma = proc_conf.use_lemmatization and (lemmatizer is not None)
    
    # defines set of "edge junk" to strip (quotes, brackets, dashes)
    # to catch "word" or (word) or -word-
    _strip_chars = string.punctuation + "“”‘’–—" 
    
    local_global_counts = Counter()
    local_global_bigrams = Counter() if proc_conf.compute_bigrams else Counter()
    
    batch_accum = Counter()
    batch_rows = 0
    row_count = 0
    
    # pre-caching lemmatizer methods for speed loop
    lemmatize = lemmatizer.lemmatize if _lemma else None

    for (raw_text, date_val, cat_val) in rows_iter:
        row_count += 1
        
        # entities (Before lowercase)
        if raw_text:
            entities = extract_entities_regex(raw_text, _stopwords)
            scanner.update_entities(entities)

        # cleaning
        text = apply_text_cleaning(raw_text, clean_conf)
        
        # tokenization & filter
        filtered_tokens_line: List[str] = []
        for t in text.split():
            # 1. Internal Translation (removes internal chars like don't -> dont)
            t = t.translate(_trans)
            
            # 2. Edge Stripping (The Fix for "between and operations")
            t = t.strip(_strip_chars)
            
            if not t: continue
            if _drop_int and t.isdigit(): continue
            if len(t) < _min_len: continue
            
            # lemmatize?
            if _lemma:
                t = lemmatize(t, pos='v')
                t = lemmatize(t, pos='n')
            
            if t in _stopwords: continue
            filtered_tokens_line.append(t)
        
        if filtered_tokens_line:
            # Stats Update
            line_counts = Counter(filtered_tokens_line)
            local_global_counts.update(filtered_tokens_line)
            
            if proc_conf.compute_bigrams and len(filtered_tokens_line) > 1:
                local_global_bigrams.update(pairwise(filtered_tokens_line))
            
            # metadata update
            clean_date = clean_date_str(date_val)
            clean_cat = str(cat_val).strip() if cat_val else None
            scanner.update_metadata_stats(clean_date, clean_cat, filtered_tokens_line)

            # topic modeling batching
            batch_accum.update(line_counts)
            batch_rows += 1
            if batch_rows >= scanner.DOC_BATCH_SIZE:
                scanner.add_topic_sample(batch_accum)
                batch_accum = Counter()
                batch_rows = 0

        if progress_cb and (row_count % 2000 == 0): progress_cb(row_count)

    # flushing last batch
    if batch_accum and batch_rows > 0:
        scanner.add_topic_sample(batch_accum)

    scanner.update_global_stats(local_global_counts, local_global_bigrams, row_count)
    if progress_cb: progress_cb(row_count)
    gc.collect()

def perform_refinery_job(file_obj, chunk_size, clean_conf: CleaningConfig):
    with tempfile.TemporaryDirectory() as temp_dir:
        original_name = os.path.splitext(file_obj.name)[0]
        status_container = st.status(f"⚙️ Refining {file_obj.name}...", expanded=True)
        part_num = 1
        created_files = []
        
        try:
            file_obj.seek(0)
            df_iterator = pd.read_csv(file_obj, chunksize=chunk_size, on_bad_lines='skip', dtype=str)
            
            for chunk in df_iterator:
                for col in chunk.columns:
                    chunk[col] = chunk[col].fillna("")
                    chunk[col] = chunk[col].apply(lambda x: apply_text_cleaning(x, clean_conf))
                
                new_filename = f"{original_name}_cleaned_part_{part_num}.csv"
                temp_path = os.path.join(temp_dir, new_filename)
                chunk.to_csv(temp_path, index=False)
                created_files.append(temp_path)
                status_container.write(f"✅ Processed chunk {part_num} ({len(chunk)} rows)")
                part_num += 1
            
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                for file_path in created_files:
                    zip_file.write(file_path, arcname=os.path.basename(file_path))
            
            zip_buffer.seek(0)
            status_container.update(label="🎉 Refinery Job Complete!", state="complete", expanded=False)
            return zip_buffer
            
        except Exception as e:
            status_container.update(label="❌ Error", state="error")
            st.error(f"Refinery Error: {str(e)}")
            return None


# 📊 UI, analytics renderers
# ==========================================

def calculate_text_stats(counts: Counter, total_rows: int) -> Dict:
    total_tokens = sum(counts.values())
    unique_tokens = len(counts)
    avg_len = sum(len(word) * count for word, count in counts.items()) / total_tokens if total_tokens else 0
    return {
        "Total Rows": total_rows, "Total Tokens": total_tokens,
        "Unique Vocabulary": unique_tokens, "Avg Word Length": round(avg_len, 2),
        "Lexical Diversity": round(unique_tokens / total_tokens, 4) if total_tokens else 0
    }

def calculate_npmi(bigram_counts: Counter, unigram_counts: Counter, total_words: int, min_freq: int = 3) -> pd.DataFrame:
    results = []
    if not bigram_counts: return pd.DataFrame(columns=["Bigram", "Count", "NPMI"])
    epsilon = 1e-10 
    for (w1, w2), freq in bigram_counts.items():
        if freq < min_freq: continue
        count_w1 = unigram_counts.get(w1, 0)
        count_w2 = unigram_counts.get(w2, 0)
        if count_w1 == 0 or count_w2 == 0: continue
        prob_bigram = freq / total_words
        try: 
            pmi = math.log(prob_bigram / ((count_w1 / total_words) * (count_w2 / total_words)))
        except ValueError: continue
        
        log_prob_bigram = math.log(prob_bigram)
        if abs(log_prob_bigram) < epsilon: npmi = 1.0
        else: npmi = pmi / -log_prob_bigram
        results.append({"Bigram": f"{w1} {w2}", "Count": freq, "NPMI": round(npmi, 3)})
        
    df = pd.DataFrame(results)
    if df.empty: return pd.DataFrame(columns=["Bigram", "Count", "NPMI"])
    return df.sort_values("NPMI", ascending=False)

def calculate_tfidf(scanner: StreamScanner, top_n=50) -> pd.DataFrame:
    # IDF = log(total docs / doc freq)
    # TF (Global) = total count / total words (simplified for sketch)
    total_docs = len(scanner.topic_docs)
    if total_docs == 0: return pd.DataFrame()
    
    results = []
    # Analyze only terms that appear in at least 2 docs to avoid noise
    candidates = [t for t, c in scanner.doc_freqs.items() if c > 1]
    
    for term in candidates:
        tf = scanner.global_counts[term]
        df = scanner.doc_freqs[term]
        idf = math.log(total_docs / (1 + df))
        score = tf * idf
        results.append({"Term": term, "TF (Count)": tf, "DF (Docs)": df, "Keyphrase Score": round(score, 2)})
        
    df = pd.DataFrame(results)
    if df.empty: return df
    return df.sort_values("Keyphrase Score", ascending=False).head(top_n)

def perform_topic_modeling(synthetic_docs: List[Counter], n_topics: int, model_type: str) -> Optional[List[Dict]]:
    if not DictVectorizer or len(synthetic_docs) < 1: return None
    vectorizer = DictVectorizer(sparse=True)
    dtm = vectorizer.fit_transform(synthetic_docs)
    n_samples, n_features = dtm.shape
    if n_samples == 0 or n_features == 0: return None
    
    safe_n_topics = min(n_topics, min(n_samples, n_features)) if model_type == "NMF" else min(n_topics, n_samples)
    if safe_n_topics < 1: return None

    model = None
    try:
        if model_type == "LDA": model = LatentDirichletAllocation(n_components=safe_n_topics, random_state=42, max_iter=10)
        elif model_type == "NMF": model = NMF(n_components=safe_n_topics, random_state=42, init='nndsvd')
        model.fit(dtm)
    except ValueError: return None
    
    feature_names = vectorizer.get_feature_names_out()
    topics = []
    for topic_idx, topic in enumerate(model.components_):
        top_indices = topic.argsort()[:-7:-1]
        top_words = [feature_names[i] for i in top_indices]
        strength = sum(topic[i] for i in top_indices)
        topics.append({"id": topic_idx + 1, "words": top_words, "strength": strength})
    return topics

def perform_bayesian_sentiment_analysis(counts: Counter, sentiments: Dict[str, float], pos_thresh: float, neg_thresh: float) -> Optional[Dict]:
    if not beta_dist: return None
    pos_count = sum(counts[w] for w, s in sentiments.items() if s >= pos_thresh)
    neg_count = sum(counts[w] for w, s in sentiments.items() if s <= neg_thresh)
    total_informative = pos_count + neg_count
    if total_informative < 1: return None

    alpha_post = 1 + pos_count
    beta_post = 1 + neg_count
    mean_prob = alpha_post / (alpha_post + beta_post)
    lower_ci, upper_ci = beta_dist.ppf([0.025, 0.975], alpha_post, beta_post)
    x = np.linspace(0, 1, 300)
    y = beta_dist.pdf(x, alpha_post, beta_post)
    return {
        "pos_count": pos_count, "neg_count": neg_count, "total": total_informative,
        "mean_prob": mean_prob, "ci_low": lower_ci, "ci_high": upper_ci,
        "x_axis": x, "pdf_y": y
    }

def render_workflow_guide():
    with st.expander("📘 Comprehensive App Guide: How to use this Tool", expanded=False):
        st.markdown("""
        ### 🌟 What is Signal Foundry?
        Signal Foundry turns messy text into structured signal. Feed it PDFs, CSVs, transcripts, PowerPoints, pasted notes, or pre-computed sketches and it will surface the recurring language, relationships, themes, entities, and maturity signals hiding inside the corpus.

        It is designed for fast exploratory analysis:
        *   **What is this corpus about?**
        *   **What terms actually matter?**
        *   **Who keeps showing up?**
        *   **What concepts cluster together?**
        *   **What level of maturity or sophistication does the language signal?**

        ---

        ### 🚀 How to Use It (Three Paths)

        **Path A – Instant Scan (best default for most users)**  
        1. Upload one or more files in the sidebar, or paste URLs / raw text.  
        2. Leave the default cleaning settings on unless you know you need something different.  
        3. If you uploaded multiple items, click **"⚡ Scan ALL Items"**.  
        4. Read the Word Cloud, Keyphrases, Entities, and Graph first.  
        5. Move into Trends, Topic Modeling, Sentiment, or Maturity only after the first pass makes sense.

        **Path B – Surgical Precision (for structured files)**  
        Use this when columns matter.
        1. Expand the **Config** panel for a CSV or Excel file.
        2. Choose the text columns you want analyzed.
        3. Optionally choose a date column for trend analysis.
        4. Optionally choose a category column for future slicing / grouping work.
        5. Run the file individually with **Start Scan**.

        **Path C – Offline / Secure-Site Mode**  
        Run `harvester.py` on your secure server, then upload the resulting `.json` sketch.
        **Result:** You still get the visualization dashboard without moving raw source text into the Streamlit app.

        ---

        ### 🧭 Recommended Reading Order

        If you are new to the app, read the outputs in this order:
        1. **Word Cloud + Stats** for the big picture
        2. **Keyphrases / NPMI** for specific language and terminology
        3. **Entities** for people, orgs, systems, or named things
        4. **Graph** for relationships and clusters
        5. **Trends** if dates were supplied
        6. **Topic Modeling** when you want automatic thematic buckets
        7. **Maturity** when you want an interpretive framework rather than raw term counts

        ---

        ### 🧠 How to Read the Output

        *   **Word Cloud + Stats** → Fast orientation. Use this to confirm you scanned the right thing.
        *   **Keyphrases (TF-IDF)** → What is unusually specific to this corpus, not just frequent.
        *   **Sticky Concepts (NPMI)** → Which word pairs behave like real concepts rather than accidental neighbors.
        *   **Entities** → Which named people, organizations, programs, or systems dominate the discussion.
        *   **Network Graph** → Which terms cluster together and where the semantic gravity sits.
        *   **Topic Modeling** → Machine-generated thematic buckets when you want a rough thematic map.
        *   **Bayesian Sentiment** → Direction plus confidence, useful only when there is enough data.
        *   **Maturity Model** → An interpretive scoring layer that reads the language as capability signal.

        ---

        ### 🧪 Topic Modeling: LDA vs NMF

        **Choose LDA when:**
        *   You are scanning longer documents, reports, policies, or dense PDFs.
        *   You expect each chunk to contain a mix of themes.
        *   You want broader, more blended topic buckets.

        **Choose NMF when:**
        *   You are scanning shorter rows, tickets, chats, transcripts, or survey comments.
        *   You want sharper, more distinct topic buckets.
        *   Your input feels repetitive but specific.

        **If the topics look bad:**
        *   Increase **Rows per Doc** for long reports and books.
        *   Decrease **Rows per Doc** for chats, support logs, and transcripts.
        *   Add boring boilerplate terms to **Stopwords**.
        *   Try the other model. Topic modeling is exploratory, not definitive.

        ---

        ### ⚡ Practical Tips

        *   **Additive Analysis:** Leave **Clear previous data** unchecked only when you intentionally want to merge new material into the current corpus.
        *   **Graph is a blob?** → Raise **Min Link Frequency**.  
        *   **Graph is empty islands?** → Lower it.  
        *   **Seeing garbage words?** → Add them to **Stopwords** box.  
        *   **Seeing "run" and "running"?** → Turn on **Lemmatization**.  
        *   **Dates not showing up in Trends?** → Make sure you selected the date column before scanning.
        *   **Need a shareable verification artifact?** → Download the **Hybrid Signature** (QR + Heatmap).

        **Bottom line:** start simple, confirm the scan makes sense, then layer on the more interpretive tools.
        """)

def render_lit_case_study():
    # We use Unicode "Math Sans" characters to simulate bold/italics in the title
    # Italic 'another': 𝘢𝘯𝘰𝘵𝘩𝘦𝘳
    # Bold 'specific': 𝘀𝗽𝗲𝗰𝗶𝗳𝗶𝗰
    title = "🔦 Spotlight: Digital Humanities & Literary Forensics (𝘢𝘯𝘰𝘵𝘩𝘦𝘳 𝘀𝗽𝗲𝗰𝗶𝗳𝗶𝗰 Case Study)"
    
    with st.expander(title, expanded=False):
        st.markdown("""
        ### The Scenario
        **The Artifact:** The full text of Ovid's **<a href="https://www.gutenberg.org/files/21765/21765-h/21765-h.htm" target="_blank">"Metamorphoses"</a>** (via Project Gutenberg URL).
        **The User:** A Digital Humanities Researcher or Student.
        **The Goal:** To rapidly map the "Pantheon" of characters and distinguish the original narrative from the translator's artifacts.

        ---

        ### 1. The "Pantheon Map" (Entities Tab)
        *   **The Question:** "Who are the dominant power players in this 15-book epic?"
        *   **The Signal:** Capitalized Name Extraction.
        *   **The Result:** The engine immediately surfaces **"Jupiter," "Apollo," "Ceres,"** and **"Minerva"** as the top nodes.
        *   **The Value / Insight:** Without reading a single line, you have a hierarchical map of the Roman deities driving the plot.

        ### 2. The "Translator's Fingerprint" (NPMI & Bigrams)
        *   **The Question:** "Is this pure text, or is there structural noise?"
        *   **The Signal:** Sticky Concepts (Bigrams).
        *   **The Result:** The engine identifies **"Clarke translates"** and **"-ver Clarke"** as top phrases.
        *   **The Value / Insight:** **Forensic Separation.** The engine detected that *John Clarke* (the translator) is statistically inseparable from the text. It highlights "Data Hygiene" issues—showing you exactly what "boilerplate" needs to be cleaned (e.g., "Project Gutenberg" headers) before deep analysis.

        ### 3. The "Narrative Arcs" (Topic Modeling)
        *   **The Question:** "What are the distinct recurring themes?"
        *   **The Signal:** NMF/LDA Mathematical Bucketing.
        *   **The Result:**
            *   **Topic A:** [Daughter, Jupiter, Cadmus, Wife] -> *The Genealogy & Origin Myths.*
            *   **Topic B:** [Thou, Thee, Thus, Said] -> *The Dialogue & Poetic Structure.*
        *   **The Value / Insight:** The engine successfully separates the *Style* (Archaic English) from the *Substance* (Mythological Events).

        ### 4. The "Semantic Network" (Graph Tab)
        *   **The Question:** "How do the main characters interact?"
        *   **The Signal:** Proximity-based linking.
        *   **The Result:** "jupiter" is the central "hub" node, with spokes connecting to various "nymphs" and "daughters."
        *   **The Value / Insight:** Visualizes the centralized power structure of the mythology, confirming Jupiter as the primary driver of the transformations.
        """, unsafe_allow_html=True)

def render_auto_insights(scanner, proc_conf):
    # Only run if we have data
    if not scanner.global_counts: return

    # --- 1. PREPARE DATA ---
    # Entities
    top_ents = scanner.entity_counts.most_common(3)
    ent_str = ", ".join([f"**{e[0]}**" for e in top_ents]) if top_ents else "(No entities detected)"
    
    # Sticky Concepts (NPMI)
    df_npmi = calculate_npmi(scanner.global_bigrams, scanner.global_counts, scanner.total_rows_processed)
    top_npmi = df_npmi.head(3)["Bigram"].tolist() if not df_npmi.empty else []
    npmi_str = ", ".join([f"**{b}**" for b in top_npmi]) if top_npmi else "(No strong phrases found)"
    
    # Tech Signal (TF-IDF)
    df_tfidf = calculate_tfidf(scanner, 20)
    top_idf = df_tfidf.head(3)["Term"].tolist() if not df_tfidf.empty else []
    idf_str = ", ".join([f"**{t}**" for t in top_idf]) if top_idf else "(Not enough documents for TF-IDF)"

    # -render reporting
    with st.expander("⚡ High-Level Signal Report (Auto-Generated, ymmv)", expanded=True):
        st.markdown(f"""
        ### 1. The "Stakeholder Map" (Entities)
        *   **The Question:** "Who are the dominant actors or organizations?"
        *   **The Signal:** Capitalized Name Extraction.
        *   **The Result:** {ent_str}
        *   **The Insight:** These nodes appear most frequently, suggesting they are the primary drivers of the narrative or the key subjects of the file.

        ### 2. The "Sticky Concepts" (Phrase Significance)
        *   **The Question:** "What is the specific 'Term of Art' or jargon here?"
        *   **The Signal:** NPMI (Normalized Pointwise Mutual Information).
        *   **The Result:** {npmi_str}
        *   **The Insight:** These words appear together mathematically more often than random chance, indicating they represent specific concepts (e.g. "Credit Card" vs "Red Card") rather than generic language.

        ### 3. The "Technical Signal" (Keyphrases)
        *   **The Question:** "What makes this specific document unique?"
        *   **The Signal:** TF-IDF (Inverse Document Frequency).
        *   **The Result:** {idf_str}
        *   **The Insight:** While words like "the" or "report" might be frequent, *these* specific words are statistically unique to this dataset, representing its core technical signature.
        """)

def render_neurotech_case_study():
    with st.expander("🔦 Spotlight: Analyzing Mi|itary Neurotechno|ogy (a very *specific* Case Study)", expanded=False):
        st.markdown("""
        ### The Scenario
        **The Artifact:** A dense, 50-page UNIDIR report titled <a href="https://unidir.org/wp-content/uploads/2025/11/UNIDIR_Neurotechnology_Military-Domain_A-Primer.pdf" target="_blank"><b>"Neurotechnology in the Military Domain"</b></a>.
        **The User:** A Defense Ana|yst with 5 minutes to extract actionable inte||igence.
        **The Goal:** Move beyond "what is this paper about?" to "what are the threats and opportunities?"

        ---

        ### 1. The "Sticky Concepts" (NPMI Tab)
        *   **The Question:** "What specific types of risks are discussed?"
        *   **The Signal:** The engine finds words that mathematically *stick together* more than random chance.
        *   **The Result:** It surfaces **"Dua| Use"** and **"Cognitive Liberty."**
        *   **The Insight:** The strategic risk isn't just new weap0ns; it is civi|ian medica| techno|ogy being repurposed for mi|itary app|ications (Dua| Use), necessitating a legal/ethical framework (Liberty).

        ### 2. The "Technical Signal" (Keyphrases Tab)
        *   **The Question:** "Do I need to worry about brain implants yet?"
        *   **The Signal:** TF-IDF filters out generic words to find unique technical terms.
        *   **The Result:** High scores for **"Non-invasive,"** **"Transcranial,"** and **"Wearable."**
        *   **The Insight:** The immediate operationa| reality is external headsets/helmets, not surgical implants.

        ### 3. The "Semantic Network" (Graph Tab)
        *   **The Question:** "How is the techno|ogy being applied?"
        *   **The Signal:** The Graph links words based on proximity in the text.
        *   **The Result:** 
            *   Cluster A links **"Stimulation"** to **"Performance"** (Enhancement/Super-Soldiers).
            *   Cluster B links **"Stimulation"** to **"Interrogation"** (Weap0nization/T0rture).
        *   **The Insight:** The paper treats "Enhancement" and "Weaponization" as distinct operational clusters.

        ### 4. The "Stakeholder Map" (Entities Tab)
        *   **The Question:** "Who is involved?"
        *   **The Signal:** Capitalized Name Extraction.
        *   **The Result:** **"D@RPA," "Neura|ink," "Geneva Convention," "Human Rights Council."**
        *   **The Insight:** Identifies the funding sources (DARP@) vs. the regulatory blockers (Geneva).
        """, unsafe_allow_html=True)

#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
def render_maturity_guide():
    with st.expander("🏆 Guide: Understanding the Maturity Models", expanded=False):
        st.markdown("""
        ### What is the Maturity Engine?

        While standard analytics count *what* words appear, the Maturity Engine measures the **intent and capability** behind those words. It compares your text against known frameworks of organizational development.

        ---

        ### 🎭 The Personas (Context Matters!)

        Words change meaning based on context. You must select the model that matches your document's intent.

        #### 1. 🏫 EdTech & LMS Ops (5-Level)
        *Best for: course designs, syllabi, LMS usage reviews, instructional strategy documents.*
        *   **L1 (Repository):** Static content dumps and basic file hosting.
        *   **L3 (Integrated):** Connected systems, tools, APIs, and interoperability.
        *   **L5 (Strategic):** Institutional thinking around equity, accessibility, and transformation.

        #### 2. 🏢 General Business Ops (5-Level, CMMI)
        *Best for: project updates, business operations notes, internal reports, emails.*
        *   **L1 (Reactive):** Urgent, manual, firefighting language.
        *   **L4 (Measured):** KPI, ROI, forecasting, and quantitative management language.

        #### 3. ⚖️ Policy & Governance (5-Level)
        *Best for: policy documents, compliance reviews, NGO/government publications.*
        *   **L1 (Enforcement):** Violations, bans, restrictions, sanctions.
        *   **L5 (Systemic):** Sustainable, resilient, holistic, ecosystem-level language.

        #### 4. 🎓 Brightspace Maturity Model — 12-Domain (3-Tier)
        *Best for: meeting transcripts, coaching notes, strategic reviews, TAM/client discussions.*

        This persona evaluates the client across **12 separate domains**, each scored on 3 tiers:

        | Tier | Label | Score Range | Meaning |
        |------|-------|-------------|---------|
        | 🔴 | **Foundational** | 1.0–1.49 | Basic, manual, reactive |
        | 🟠 | **Advanced** | 1.50–2.49 | Structured, proactive, repeatable |
        | 🟢 | **Leading Edge** | 2.50–3.0 | Strategic, scalable, innovation-oriented |

        **The 12 Domains**
        1. Platform & Technical Administration
        2. Curriculum Development & Delivery
        3. Student Engagement & Success
        4. Data & Learning Analytics
        5. Assessment & Evaluation
        6. Instructor Efficiency
        7. Change Management
        8. Knowledge & Resource Management
        9. Accessibility & Compliance
        10. User Support & Training
        11. Innovation & Emerging Technologies
        12. Collaboration & Communication

        ---

        ### 🧮 How Scoring Works

        **For the 5-level personas**
        *   Weighted average across the detected maturity-level signals.
        *   Score range: **1.0 to 5.0**

        **For the 12-domain Brightspace persona**
        *   Each domain is scored independently using tiers 1-3.
        *   The composite score is the average of all domains with at least one signal.
        *   Score range: **1.0 to 3.0**
        *   Domains with zero signal are shown as **No Data** and excluded from the composite.

        **Phrase weighting**
        *   Multi-word phrases receive extra weight relative to single words because they are usually more specific signals.

        ---

        ### ⚙️ Important Processing Note

        Maturity scoring uses the same cleaned token stream as the rest of the app. That means settings like:
        *   **Min Word Len**
        *   **Stopwords**
        *   **Remove Generic Filler Words and Prepositions**
        *   **Use Lemmatization**
        *   **Keep Hyphens**
        *   **Chat / HTML / URL cleanup**

        can all affect the maturity result.

        ---

        ### 🕸️ How to Read the Charts

        **5-Level Radar**
        *   Round shape = balanced maturity
        *   Sharp spikes = uneven capability

        **12-Domain Radar**
        *   One spoke per domain
        *   Outer edge = stronger maturity signal
        *   Uneven shape = domain-specific strengths and gaps

        **12-Domain Breakdown Bars**
        *   Red = Foundational share
        *   Orange = Advanced share
        *   Green = Leading Edge share

        ---

        ### ⚠️ Caveats

        > **Signal vs. reality:** The maturity engine measures language, not actual implementation. Treat it as a conversation starter and evidence aid, not a final audit.

        > **Persona mismatch:** If you scan a generic business strategy memo with the Brightspace model, the score may look plausible but be contextually wrong.

        > **Data quality matters:** Transcripts, coaching notes, and strategy discussions usually work better than exported system logs or boilerplate-heavy files.
        """)
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<

def render_use_cases():
    with st.expander("📖 Playbook: High-Value Use Cases", expanded=False):
        st.markdown("""
        ### 🎓 Specialized: Education & EdTech
        
        #### 1. The "Syllabus Audit" (Maturity Scan)
        *   **Goal:** Determine if a department is truly "Modernizing" or just digitizing old habits.
        *   **Action:** Bulk scan 50 syllabi/course descriptions using the **EdTech Persona**.
        *   **Signal:**
            *   High **L1 (Repository)** scores = The LMS is just a file dump.
            *   High **L3 (Integrated)** scores = They are actually using the ecosystem tools.
        
        #### 2. Vendor RFP Analysis
        *   **Goal:** Cut through sales fluff.
        *   **Action:** Scan a vendor's whitepaper or proposal.
        *   **Signal:** Do they speak **L2 (Features/Tools)** or **L5 (Partnership/Success)**? A strategic partner should score >3.5.
        
        ---

        ### 🏢 General: Corporate & Intelligence
        
        #### 3. Strategic Alignment Check
        *   **Goal:** See if the IT Department aligns with the C-Suite.
        *   **Action:** Scan the CTO's emails vs. the CEO's annual letter.
        *   **Signal:** If the CEO is L5 (Vision) and IT is L1 (Firefighting), you have an execution gap.

        #### 4. "Chain of Custody" Verification
        *   **Goal:** Prove that a visualization came from *this* specific file.
        *   **Action:** Generate the **Hybrid Signature** (QR Code) in the Graph tab.
        *   **Value:** Encrypts the document hash into the image. If the screenshot leaks, you can prove exactly which source file generated it.

        #### 5. Crisis Timeline Reconstruction
        *   **Goal:** Pinpoint when a project went off the rails.
        *   **Action:** Scan weekly status reports with the **Time-Travel Slider**.
        *   **Signal:** Watch for the crossover point where **"Plan" (L2)** words drop and **"Fix/Urgent" (L1)** words spike.

        ---

        ### 🎓 Specialized: TAM & Client Maturity Assessment

        #### 6. The "Client Maturity Snapshot" (12-Domain Scan)
        *   **Goal:** Quickly gauge where a client stands across all 12 Brightspace admin maturity domains.
        *   **Action:** Upload 3–5 meeting transcripts or coaching session recordings (VTT/PDF) from a single client. Select the **🎓 Brightspace Maturity Model** persona.
        *   **Signal:**
            *   The **Radar Chart** instantly shows which domains the client discusses at a strategic level vs. which are still basic.
            *   The **Breakdown Chart** reveals the distribution — is the client *mostly* Foundational with a few Advanced pockets, or broadly Advanced with Leading Edge potential?
        *   **Action:** Use the domain cards to identify the 2–3 domains with the lowest scores and build your coaching plan around them.

        #### 7. Client Progress Over Time (Longitudinal Tracking)
        *   **Goal:** Measure whether a client's maturity is improving after coaching.
        *   **Action:** Export a snapshot after each major engagement (quarterly reviews, coaching milestones). Later, upload all snapshots to the Progress Tracking section.
        *   **Signal:** Compare the composite scores and per-domain scores across dates. Are the reds turning orange? Are the oranges turning green?
        *   **Value:** Concrete, data-backed evidence of TAM impact for quarterly reviews and stakeholder reporting.

        #### 8. Cross-Client Benchmarking
        *   **Goal:** Identify which clients need the most attention.
        *   **Action:** Scan transcripts from 3–4 different clients separately. Note each client's composite score and lowest-scoring domains.
        *   **Signal:** Client A scores 2.4 composite, Client B scores 1.3. Client B needs more coaching. Client A might be ready for Leading Edge challenges.
        *   **Value:** Helps TAMs and OS leadership prioritize resource allocation across the client portfolio.
        """)

#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
def render_analyst_help():
    with st.expander("🎓 Analyst's Guide & Troubleshooting", expanded=False):
        st.markdown("""
        ### Quick Diagnostic Rule

        If the **Word Cloud** and **Frequency Tables** look wrong, the downstream outputs will also be wrong. Fix cleaning, stopwords, file columns, or scan mode first.

        ---

        **Symptom: The app output feels wrong right away**
        * **Fix:** Confirm you scanned the correct file, sheet, and text column.
        * **Fix:** Check whether **Clear previous data** was off and old data was accidentally merged into the current scan.
        * **Fix:** Look at the top frequency table. If the top words are boilerplate, add them to **Stopwords** and rescan.

        **Symptom: Important acronyms disappeared**
        * **Likely cause:** **Min Word Len** is too high.
        * **Fix:** Lower **Min Word Len** to 2 or 3 for acronym-heavy corpora such as API, LTI, SSO, AI, SIS, or WCAG.

        **Symptom: Maturity results feel too low or missing domains**
        * **Fix:** Confirm the selected maturity persona matches the source material.
        * **Fix:** Lower **Min Word Len** if important short terms are being filtered out.
        * **Fix:** Keep **Bigrams** enabled so phrase signals like "student success" and "change management" can be detected.
        * **Fix:** Feed the model more relevant meeting notes, transcripts, coaching notes, or strategic review material.
        * **Interpretation:** A **No Data** domain means no mapped language was detected for that domain. It does not prove the client lacks that capability.

        **Symptom: Maturity result feels too generic**
        * **Fix:** Use the domain detail cards. The top linguistic drivers tell you exactly why the score was produced.
        * **Fix:** If drivers look like boilerplate, add that boilerplate to **Stopwords** and rescan.
        * **Fix:** If a meaningful client-specific phrase is missing, the maturity vocabulary may need expansion.

        **Symptom: Topics look like gibberish or random words**
        * **Fix:** Check **Rows per Doc**.
            * For chats, tickets, comments, and transcripts, use 1–5.
            * For reports, PDFs, books, and long documents, use 100+.
        * **Fix:** Try **NMF** for shorter, cleaner records.
        * **Fix:** Try **LDA** for longer mixed-content documents.
        * **Fix:** Add recurring boilerplate terms to **Stopwords**.

        **Symptom: The Network Graph is a giant blob**
        * **Fix:** Increase **Min Link Frequency** to cut weak connections.
        * **Fix:** Increase **Repulsion** to push nodes apart.
        * **Fix:** Lower **Max Nodes** to focus on the most important terms.

        **Symptom: The Graph has disconnected islands**
        * **Fix:** Decrease **Min Link Frequency** to reveal subtler connections.
        * **Fix:** Increase **Edge Length** to give clusters room to breathe.

        **Symptom: Seeing duplicates such as "run" and "running"**
        * **Fix:** Enable **Use Lemmatization**. This can merge word variations into a shared root form.

        **Symptom: High-ranking words are boring, such as "page", "copyright", or "transcript"**
        * **Fix:** These are corpus artifacts. Add them to **Stopwords** and rescan.
        * **Fix:** For transcripts, keep **Remove Chat Artifacts** enabled.

        **Symptom: Trend charts are empty**
        * **Fix:** Re-scan after selecting the correct date column in the file config panel.
        * **Fix:** Check whether source dates are parseable, such as 2025-01-31, 01/31/2025, or Jan 31 2025.

        **Symptom: AI Analyst answer seems too vague**
        * **Fix:** The AI Analyst only sees the statistical sketch, not raw documents.
        * **Fix:** Ask narrower questions about visible outputs, such as "Which maturity domains look weakest?" or "Which graph clusters appear most central?"
        """)
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<
# visualization helpers
@st.cache_data(show_spinner="Analyzing term sentiment...")
def get_sentiments(_analyzer, terms: Tuple[str, ...]) -> Dict[str, float]:
    if not _analyzer or not terms: return {}
    return {term: _analyzer.polarity_scores(term)['compound'] for term in terms}

def create_sentiment_color_func(sentiments: Dict[str, float], pos_color, neg_color, neu_color, pos_thresh, neg_thresh):
    def color_func(word, font_size, position, orientation, random_state=None, **kwargs):
        score = sentiments.get(word, 0.0)
        if score >= pos_thresh: return pos_color
        elif score <= neg_thresh: return neg_color
        else: return neu_color
    return color_func

def get_sentiment_category(score: float, pos_threshold: float, neg_threshold: float) -> str:
    if score >= pos_threshold: return "Positive"
    if score <= neg_threshold: return "Negative"
    return "Neutral"

def build_wordcloud_figure_from_counts(counts, max_words, width, height, bg_color, colormap, font_path, random_state, color_func):
    limited = dict(counts.most_common(max_words))
    if not limited: return plt.figure(), None
    wc = WordCloud(width=width, height=height, background_color=bg_color, colormap=colormap, font_path=font_path, random_state=random_state, color_func=color_func, collocations=False, normalize_plurals=False).generate_from_frequencies(limited)
    fig, ax = plt.subplots(figsize=(max(6, width/100), max(3, height/100)), dpi=100)
    ax.imshow(wc, interpolation="bilinear"); ax.axis("off"); plt.tight_layout()
    return fig, wc

def fig_to_png_bytes(fig):
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", pad_inches=0.1)
    buf.seek(0)
    return buf

# 🤖 AI logic
def call_llm_and_track_cost(system_prompt: str, user_prompt: str, config: dict):
    try:
        client = openai.OpenAI(api_key=config['api_key'], base_url=config['base_url'])
        response = client.chat.completions.create(
            model=config['model_name'],
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
        )
        in_tok = 0
        out_tok = 0
        if hasattr(response, 'usage') and response.usage:
            in_tok = response.usage.prompt_tokens
            out_tok = response.usage.completion_tokens
        
        cost = (in_tok * config['price_in'] / 1_000_000) + (out_tok * config['price_out'] / 1_000_000)
        st.session_state['total_tokens'] += (in_tok + out_tok)
        st.session_state['total_cost'] += cost
        return response.choices[0].message.content
    except Exception as e:
        return f"AI Error: {str(e)}"


# 🚀 main app ui
# ==========================================

st.set_page_config(page_title="Signal Foundry", layout="wide")
st.toast("app loaded/updated successfully", icon="🚀") # cache buster
st.title("🧠 Signal Foundry: Unstructured Data Analytics")
st.markdown("### *(or: data geiger counter~)*")

# Initialize NLP globally
analyzer, lemmatizer = setup_nlp_resources()

# --- SIDEBAR (Global Inputs) ---
with st.sidebar:
    st.header("📂 Data Input")
    uploaded_files = st.file_uploader(
        "Upload Files",
        type=["csv", "xlsx", "vtt", "txt", "json", "pdf", "pptx"],
        accept_multiple_files=True,
        help="Upload one or more source files. CSV/XLSX are best for structured data, VTT/TXT for transcripts, PDF/PPTX for document decks, and JSON for offline sketches.",
    )
    
    # --modified logic (Additive Mode Safety)
    
    # checking if there's currently data
    has_data = st.session_state['sketch'].total_rows_processed > 0
    
    # showing the checkbox
    clear_on_scan = st.checkbox(
        "Clear previous data",
        value=False,
        help="Turn this on when you want a fresh analysis. Leave it off only when you intentionally want to add new files into the current corpus.",
    )
    
    # 'banner' logic: only showing if there's ambiguity
    if has_data and not clear_on_scan:
        st.info("⚠️ **Additive Mode Active:** New scans will be ADDED to current results. Check the 'Clear previous data' check-box above to start fresh", icon="ℹ️")
    elif has_data and clear_on_scan:
        st.caption("✅ Next scan will overwrite current data.")
        
    if st.button("🗑️ Reset All"): reset_sketch(); st.rerun()
    
    # --end
    
    st.divider()

    # 1. quick imports (web/text)
    with st.expander("🌐 Quick Web / Text Paste"):
        url_input = st.text_area(
            "URLs (1 per line)",
            placeholder="https://example.com/article",
            help="Paste public URLs to scrape page text into the current scan. Best for articles, public reports, and web pages with readable body text.",
        )
        manual_input = st.text_area(
            "Manual Text Paste",
            placeholder="Paste raw text content here...",
            help="Use this for notes, excerpts, transcripts, or ad hoc text that you do not want to upload as a file.",
        )

    # 2. offline / enterprise import ("harvester")
    with st.expander("📡 Load Offline Analysis (Harvester)", expanded=False):
        st.markdown("Upload pre-computed analysis from your secure server.")
        
        # the feature gap warning
        st.caption("⚠️ **Note:** Offline sketches provide Graphs, Counts & NPMI. Time-series and Entities are disabled in this mode.")
        
        sketch_upload = st.file_uploader(
            "Upload Sketch File (.json)", 
            type=["json"],
            help="Use this for datasets >200MB. Run the 'harvester.py' script on your data server, then upload the resulting JSON here.\n\nCommand:\npython harvester.py --input data.csv --col text --output sketch.json"
        )
        
        if sketch_upload:
            file_hash = hash(sketch_upload.getvalue())
            if st.session_state.get('last_sketch_hash') != file_hash:
                try:
                    # clearing previous data to prevent corruption
                    reset_sketch()
                    
                    json_str = sketch_upload.getvalue().decode('utf-8')
                    if st.session_state['sketch'].load_from_json(json_str):
                        st.session_state['last_sketch_hash'] = file_hash
                        st.success(f"✅ Loaded Sketch: {sketch_upload.name}")
                    else:
                        st.error("❌ Invalid Schema: JSON does not match Signal Foundry structure.")
                except Exception as e:
                    st.error(f"❌ Load Error: {e}")
        
    st.divider()
   
    st.header("🔐 AI Setup", help="This AI feature reads the *metadata only*; it does not upload raw document text to the cloud provider, nor does it 'read' it.")
    
    if st.session_state['authenticated']:
        st.success("Unlocked")
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
        with st.expander("🤖 Provider Settings", expanded=True):
            ai_provider = st.radio(
                "Provider",
                ["DeepSeek AI", "xAI (Grok)", "OpenAI (GPT-4o)"],
                index=0,
            )

            if ai_provider == "DeepSeek AI":
                api_key_name = "deepseek_api_key"
                base_url = "https://api.deepseek.com"
                model_name = st.selectbox(
                    "Model",
                    ["deepseek-v4-flash", "deepseek-v4-pro"],
                    index=0,
                )
                # Pricing can change; these placeholders keep the existing
                # cost-estimator plumbing working without affecting API calls.
                price_in, price_out = 0.00, 0.00

            elif "OpenAI" in ai_provider:
                api_key_name = "openai_api_key"
                base_url = None
                model_name = st.selectbox("Model", ["gpt-4o", "gpt-4o-mini"])
                price_in, price_out = (0.15, 0.60) if "mini" in model_name else (2.50, 10.00)

            else:
                api_key_name = "xai_api_key"
                base_url = "https://api.x.ai/v1"
                model_name = "grok-4-0709"
                price_in, price_out = 3.00, 15.00

            api_key = st.secrets.get(api_key_name)
            if not api_key:
                api_key = st.text_input(f"Enter {api_key_name}", type="password")

            ai_config = {
                "api_key": api_key,
                "base_url": base_url,
                "model_name": model_name,
                "price_in": price_in,
                "price_out": price_out,
            }
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<

        with st.expander("💰 Cost Estimator", expanded=False):
            c1, c2 = st.columns(2)
            c1.markdown(f"**Tokens:**\n{st.session_state['total_tokens']:,}")
            c2.markdown(f"**Cost:**\n`${st.session_state['total_cost']:.5f}`")
            if st.button("Reset Cost"):
                st.session_state['total_cost'] = 0.0
                st.session_state['total_tokens'] = 0
                st.rerun()
        if st.button("Logout"): logout(); st.rerun()
    else:
        st.text_input("Password", type="password", key="password_input", on_change=perform_login)
        if st.session_state['auth_error']: st.error("Incorrect password")

    st.divider()
    
    # [NEW] Health Check Section
    with st.expander("❤️ System Health", expanded=False):
        st.markdown(f"{'✅' if requests else '❌'} **Web Scraper** (requests)")
        st.markdown(f"{'✅' if pypdf else '❌'} **PDF Support** (pypdf)")
        st.markdown(f"{'✅' if openpyxl else '❌'} **Excel Support** (openpyxl)")
        st.markdown(f"{'✅' if pptx else '❌'} **PowerPoint** (python-pptx)")
        st.markdown(f"{'✅' if nltk else '❌'} **Sentiment** (nltk)")
        st.markdown(f"{'✅' if LatentDirichletAllocation else '❌'} **Topic Modeling** (sklearn)")
        st.markdown(f"{'✅' if qrcode else '❌'} **QR Generator** (qrcode)")
    
    st.header("⚙️ Configuration")
    
    st.markdown("**Cleaning**")
    clean_conf = CleaningConfig(
        remove_chat=st.checkbox("Remove Chat Artifacts", True, help="Strips metadata like timestamps, usernames (e.g., <@U1234>), and system messages from logs/transcripts to focus purely on the conversation content."),
        remove_html=st.checkbox("Remove HTML", True, help="Removes HTML tags such as <div>, <br>, and other markup from exported web or LMS content."),
        remove_urls=st.checkbox("Remove URLs", True, help="Removes links and email addresses so they do not pollute the vocabulary and graph."),
        unescape=st.checkbox("Unescape HTML", True, help="Converts coded entities (e.g., &amp ; amp;, &amp ; quot;) back into readable symbols (&, \").")
    )
    
    st.markdown("**Processing**")
    use_lemma = st.checkbox("Use Lemmatization", False, help="Merges 'running' -> 'run'. Slower but cleaner.")
    if use_lemma and lemmatizer is None: st.warning("NLTK Lemmatizer not found.")
    
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
    proc_conf = ProcessingConfig(
        min_word_len=st.slider(
            "Min Word Len",
            1,
            10,
            4,
            help=(
                "Filters out short tokens before analysis. This affects Word Cloud, "
                "Keyphrases, Graphs, Topic Modeling, and Maturity scoring. Lower this "
                "for acronyms and short jargon such as LTI, API, SSO, or AI. Raise it "
                "only when short words are mostly noise."
            ),
        ),
        drop_integers=st.checkbox(
            "Drop Integers",
            True,
            help=(
                "Removes standalone numbers from the analysis. Turn this off if years, "
                "course codes, section numbers, ticket IDs, or numeric labels carry meaning."
            ),
        ),
        compute_bigrams=st.checkbox(
            "Bigrams",
            True,
            help=(
                "Tracks two-word phrases such as 'student success' or 'change management'. "
                "Required for NPMI, graph strength, and phrase-style maturity signals."
            ),
        ),
        use_lemmatization=use_lemma,
        translate_map=build_punct_translation(
            st.checkbox(
                "Keep Hyphens",
                help=(
                    "Useful when hyphenated terms matter, such as 'cross-functional' "
                    "or 'competency-based'. Hyphen handling can affect phrase matching."
                ),
            ),
            st.checkbox(
                "Keep Apostrophes",
                help=(
                    "Useful if apostrophes are meaningful in your text, but most corpora "
                    "are cleaner with them removed."
                ),
            ),
        )
    )
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<
    
    # stopwords
    user_sw = st.text_area(
        "Stopwords (comma-separated)",
        "firstname.lastname, jane doe, okay, ok, really",
        help="Add domain-specific junk words, recurring names, or boilerplate terms you want removed before analysis.",
    )
    phrases, singles = parse_user_stopwords(user_sw)
    clean_conf.phrase_pattern = build_phrase_pattern(phrases)
    stopwords = set(STOPWORDS).union(singles)
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
    if st.checkbox(
        "Remove Generic Filler Words and Prepositions",
        True,
        help="Removes common low-value linking words such as 'of', 'to', 'with', and similar filler terms that usually add noise rather than meaning.",
    ):
        stopwords.update(default_prepositions())

    stopwords = protect_maturity_vocabulary(stopwords)
    st.caption(
        "Safety note: maturity-model vocabulary is protected from stopword removal "
        "so key maturity signals are not accidentally filtered out."
    )
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<
    proc_conf.stopwords = stopwords
    
    st.markdown("### 🎨 Appearance")
    bg_color = st.color_picker("Background Color", "#000000")
    colormap = st.selectbox("Colormap", ["viridis", "plasma", "inferno", "magma", "cividis", "tab10", "Blues", "Reds", "Greys"], 0)
    top_n = st.number_input("Top Terms to Display", min_value=5, max_value=1000, value=51, help="Controls table depth and how many top terms or bigrams are shown in ranked outputs.")
    max_words = st.slider("Max Words (Cloud)", 50, 3000, 1000, 50, help="Controls how many words are eligible for the word cloud. Higher values create denser clouds, but can make them harder to read.")
    
    # font selection
    font_map, font_names = list_system_fonts(), list(list_system_fonts().keys())
    default_font_idx = 0
    desired_font = "DejaVu Sans Mono"
    if desired_font in font_names:
        default_font_idx = font_names.index(desired_font)
    combined_font_name = st.selectbox("Font", font_names or ["(default)"], 0, help="Applies to generated word clouds and some rendered visual outputs.")
    combined_font_path = font_map.get(combined_font_name) if font_names else None

    st.markdown("### 🔬 Sentiment")
    enable_sentiment = st.checkbox("Enable Sentiment", False, help="Adds sentiment scoring to top terms, bigrams, and the Bayesian sentiment view. Best used on opinion-rich text, not neutral technical exports.")
    if enable_sentiment and analyzer is None:
        st.error("NLTK not found.")
        enable_sentiment = False
    
    pos_threshold, neg_threshold, pos_color, neu_color, neg_color = 0.05, -0.05, '#2ca02c', '#808080', '#d62728'
    if enable_sentiment:
        c1, c2 = st.columns(2)
        with c1: pos_threshold = st.slider("pos threshold", 0.0, 1.0, 0.05, 0.01)
        with c2: neg_threshold = st.slider("neg threshold", -1.0, 0.0, -0.05, 0.01)
        c1, c2, c3 = st.columns(3)
        with c1: pos_color = st.color_picker("pos color", value=pos_color)
        with c2: neu_color = st.color_picker("neu color", value=neu_color)
        with c3: neg_color = st.color_picker("neg color", value=neg_color)

    # Updated with help text
    doc_granularity = st.select_slider(
        "Rows per Doc", 
        options=[1, 5, 10, 100, 500], 
        value=5,
        help="Defines how much text gets grouped into one synthetic document for topic modeling and TF-IDF-like behavior. Use 1-5 for chats, tickets, and transcripts. Use 100+ for books, reports, and long PDFs."
    )
    st.session_state['sketch'].set_batch_size(doc_granularity)
    
    current_settings_hash = f"{doc_granularity}_{proc_conf.min_word_len}"
    if 'last_settings_hash' not in st.session_state: 
        st.session_state['last_settings_hash'] = current_settings_hash

    if st.session_state['last_settings_hash'] != current_settings_hash:
        if st.session_state['sketch'].total_rows_processed > 0:
            reset_sketch()
            st.warning("⚙️ Settings changed. Data reset for consistency. Please Scan again.")
        st.session_state['last_settings_hash'] = current_settings_hash

    # Updated with help text
    topic_model_type = st.selectbox(
        "Topic Model", 
        ["LDA", "NMF"],
        help="LDA is better for longer documents where each chunk may mix several themes. NMF is better for shorter, cleaner records where topics should separate more sharply."
    )
    st.caption("Topic model rule of thumb: use **NMF** for shorter rows and transcripts; use **LDA** for longer reports and mixed-content documents.")
    
    # Updated with help text
    n_topics = st.slider(
        "Topics", 
        2, 10, 
        4,
        help="How many topic buckets to generate. Start with 3-5. If topics feel too broad, increase it slightly. If they feel repetitive or noisy, decrease it."
    )

# --- TABS LAYOUT ---
tab_work, tab_learn = st.tabs(["🚀 Workspace", "📚 Learn (How to Use & Use-Cases)"])

# 1. THE LEARNING TAB (Guides, Examples)
with tab_learn:
    render_workflow_guide()
    render_maturity_guide()  # added
    render_use_cases()
    render_neurotech_case_study()
    render_lit_case_study()
    render_analyst_help()

# 2. THE WORKSPACE TAB (Main Engine)
with tab_work:
    with st.expander("🛠️ Data Refinery (only if you need to split very large data files; **NOTE: sanitize first**)"):
        ref_file = st.file_uploader("CSV to Refine/split", type=['csv'])
        if ref_file and st.button("🚀 Run Refinery"):
            zip_data = perform_refinery_job(ref_file, 50000, clean_conf)
            if zip_data: st.download_button("Download ZIP", zip_data, "refined.zip", "application/zip")

    # --scanning phase
    all_inputs = list(uploaded_files) if uploaded_files else []
    if url_input:
        for u in url_input.split('\n'):
            if u.strip(): 
                txt = fetch_url_content(u.strip())
                if txt: 
                    all_inputs.append(VirtualFile(f"url_{hash(u)}.txt", txt))
                    time.sleep(URL_SCRAPE_RATE_LIMIT_SECONDS) # RATE LIMITING
    if manual_input: all_inputs.append(VirtualFile("manual.txt", manual_input))

    if all_inputs:
        st.subheader("🚀 Scanning Phase")
        
        # --batch scanner
        # only shows this button if we have more than 1 file/url
        if len(all_inputs) > 1:
            # dynamic button labeling
            item_count = len(all_inputs)
            has_data = st.session_state['sketch'].total_rows_processed > 0
            
            if has_data:
                if clear_on_scan:
                    # data exists + clear is TRUE -> "re-scan / overwrite"
                    batch_btn_label = f"♻️ Re-Scan ALL {item_count} Items (Overwrite)"
                    batch_btn_help = "This will WIPE the current analysis and process the file list from scratch."
                else:
                    # if data exists + clear is FALSE -> "add" (risking duplicates)
                    batch_btn_label = f"➕ Scan & Add ALL {item_count} Items (Additive)"
                    batch_btn_help = "⚠️ CAUTION: This will process these files and ADD them to the existing results. If you have already scanned them, this will create DUPLICATES. Use 'Clear previous data' to start fresh."
            else:
                # no data -> standard start
                batch_btn_label = f"⚡ Start Batch Scan ({item_count} Items)"
                batch_btn_help = "Process all items in the list above."

            if st.button(batch_btn_label, type="primary", help=batch_btn_help):
                # handles reset logic based on user checkbox
                if clear_on_scan: 
                    reset_sketch()
                
                # 2. Setup Progress
                prog_bar = st.progress(0)
                status_box = st.empty()
                
                # 3. Iterate through all files
                for i, item in enumerate(all_inputs):
                    status_box.markdown(f"**Processing {i+1}/{len(all_inputs)}:** *{item.name}*...")
                    
                    # Detect format
                    f_bytes = item.getvalue()
                    fname = item.name.lower()
                    
                    # Logic to pick the reader (Simplified detection for batch mode)
                    batch_iter = iter([])
                    if fname.endswith(".csv"):
                        headers = detect_csv_headers(f_bytes)
                        if headers:
                            # Header exists → use first column (safe default)
                            batch_iter = read_rows_csv_structured(
                                f_bytes, "auto", ",", True, [headers[0]], None, None, " "
                            )
                        else:
                            # No header detected → fall back to treating as raw text lines (current safe behavior)
                            # Alternative (more aggressive): treat as CSV with no header and join all columns
                            # Uncomment the line below if you want to capture everything in headerless CSVs
                            # batch_iter = read_rows_csv_structured(f_bytes, "auto", ",", False,
                            #                                       [f"col_{i}" for i in range(20)], None, None, " ")
                            batch_iter = read_rows_raw_lines(f_bytes)

                    elif fname.endswith(".xlsx"):
                        sheets = get_excel_sheetnames(f_bytes)
                        if sheets:
                            batch_iter = iter_excel_structured(f_bytes, sheets[0], True, ["col_0"], None, None, " ")

                    elif fname.endswith(".pdf"):
                        batch_iter = read_rows_pdf(f_bytes)

                    elif fname.endswith(".pptx"):
                        batch_iter = read_rows_pptx(f_bytes)

                    elif fname.endswith(".vtt"):
                        batch_iter = read_rows_vtt(f_bytes)

                    elif fname.endswith(".json"):
                        batch_iter = read_rows_json(f_bytes)

                    else:
                        batch_iter = read_rows_raw_lines(f_bytes)
                        
                    # Process
                    process_chunk_iter(batch_iter, clean_conf, proc_conf, st.session_state['sketch'], lemmatizer)
                    
                    # Update Progress
                    prog_bar.progress((i + 1) / len(all_inputs))
                
                status_box.success(f"✅ Batch Complete! Processed {len(all_inputs)} files.")
                st.rerun()
        # -------------------------

        for idx, f in enumerate(all_inputs):
            try:
                # resource limit check
                if f.getbuffer().nbytes > MAX_FILE_SIZE_MB * 1024 * 1024:
                    st.error(f"❌ File **{f.name}** exceeds {MAX_FILE_SIZE_MB}MB limit.")
                    continue

                file_bytes, fname, lower = f.getvalue(), f.name, f.name.lower()
                is_csv = lower.endswith(".csv")
                is_xlsx = lower.endswith((".xlsx", ".xlsm"))
                is_json = lower.endswith(".json")
                is_vtt = lower.endswith(".vtt")
                is_pdf = lower.endswith(".pdf")
                is_pptx = lower.endswith(".pptx")
                
                # Default Scan Settings
                scan_settings = {
                    "date_col": None,
                    "cat_col": None,
                    "text_cols": [],
                    "has_header": False,
                    "sheet_name": None,
                    "json_key": None
                }
                
                with st.expander(f"🧩 Config: {fname}", expanded=False):
                    if is_csv:
                        headers = detect_csv_headers(file_bytes)
                        if headers:
                            scan_settings["has_header"] = True
                            st.info(f"Detected {len(headers)} columns.")
                            scan_settings["text_cols"] = st.multiselect("Text Columns", headers, default=[headers[0]], key=f"txt_{idx}", help="Choose the column or columns that contain the main text you want analyzed.")
                            scan_settings["date_col"] = st.selectbox("Date Column (Optional)", ["(None)"] + headers, key=f"date_{idx}", help="Pick a date-like column if you want the Trends tab and time slider to work.")
                            scan_settings["cat_col"] = st.selectbox("Category Column (Optional)", ["(None)"] + headers, key=f"cat_{idx}", help="Optional grouping field. Useful when the file contains segments such as department, source, or content type.")
                            
                            if scan_settings["date_col"] == "(None)": scan_settings["date_col"] = None
                            if scan_settings["cat_col"] == "(None)": scan_settings["cat_col"] = None
                        else:
                            st.warning("No headers detected. Scanning as raw text.")
                    elif is_xlsx:
                        sheets = get_excel_sheetnames(file_bytes)
                        scan_settings["sheet_name"] = st.selectbox("Sheet", sheets, key=f"sheet_{idx}", help="Choose the Excel sheet that contains the text you want to analyze.")
                        if scan_settings["sheet_name"]:
                             scan_settings["has_header"] = st.checkbox("Has Header Row", True, key=f"xls_head_{idx}", help="Leave this on if the first row contains column names rather than actual text data.")
                    elif is_json:
                        scan_settings["json_key"] = st.text_input("JSON Key (Optional)", "", key=f"json_{idx}", help="For line-delimited JSON objects, enter the field that contains the text to analyze.")

                if st.button(f"Start Scan: {fname}", key=f"btn_{idx}"):
                    if clear_on_scan: reset_sketch()
                    bar = st.progress(0)
                    status = st.empty()
                    
                    # select iterator
                    rows_iter = iter([])
                    approx = estimate_row_count_from_bytes(file_bytes)

                    if is_csv and scan_settings["has_header"] and scan_settings["text_cols"]:
                        rows_iter = read_rows_csv_structured(
                            file_bytes, "auto", ",", True, 
                            scan_settings["text_cols"], scan_settings["date_col"], scan_settings["cat_col"], " "
                        )
                    elif is_xlsx and scan_settings["sheet_name"]:
                        rows_iter = iter_excel_structured(
                            file_bytes, scan_settings["sheet_name"], scan_settings["has_header"], 
                            ["col_0"], None, None, " " 
                        )
                    elif is_pdf:
                        rows_iter = read_rows_pdf(file_bytes)
                    elif is_pptx:
                        rows_iter = read_rows_pptx(file_bytes)
                    elif is_vtt:
                        rows_iter = read_rows_vtt(file_bytes)
                    elif is_json:
                        rows_iter = read_rows_json(file_bytes, scan_settings["json_key"])
                    else:
                        # fallback raw line reader
                        rows_iter = read_rows_raw_lines(file_bytes)
                    
                    # run
                    process_chunk_iter(rows_iter, clean_conf, proc_conf, st.session_state['sketch'], lemmatizer, lambda n: status.text(f"Rows: {n:,}"))
                    bar.progress(100)
                    status.success("Done!")
                    if not clear_on_scan: st.rerun()

            except Exception as e:
                st.error(f"Error: {e}")

    # --- analysis phase

    scanner = st.session_state['sketch']

    # [NEW] Time-Travel Logic
    counts_source = scanner.global_counts
    
    if scanner.temporal_counts:
        all_dates = sorted(scanner.temporal_counts.keys())
        if len(all_dates) > 1:
            st.divider()
            st.subheader("⏳ Time-Travel Filter")
            
            # Create the slider
            start_date, end_date = st.select_slider(
                "Filter Analysis by Timeframe",
                options=all_dates,
                value=(all_dates[0], all_dates[-1])
            )
            
            # if user moves the slider, reconstruct counts from the specific dates
            if (start_date != all_dates[0]) or (end_date != all_dates[-1]):
                subset_counts = Counter()
                # to iterate through dates in range
                in_range = False
                for d in all_dates:
                    if d == start_date: in_range = True
                    if in_range:
                        subset_counts.update(scanner.temporal_counts[d])
                    if d == end_date: break
                
                counts_source = subset_counts
                st.caption(f"Showing data from **{start_date}** to **{end_date}**")

    # now: dynamic filtering with visualization fix
    combined_counts = Counter({
        k: v for k, v in counts_source.items() 
        if len(str(k)) >= proc_conf.min_word_len
        and k not in proc_conf.stopwords # dynamic filtering
    })

    if combined_counts:
        st.divider()
        st.header("📊 Analysis Dashboard")
        
        # calculate stats upfront
        text_stats = calculate_text_stats(combined_counts, scanner.total_rows_processed)
        render_auto_insights(scanner, proc_conf)
        # main tabs
        tab_main, tab_trend, tab_ent, tab_key, tab_mat = st.tabs(["☁️ Word Cloud & Stats", "📈 Trends", "👥 Entities", "🔑 Keyphrases", "🏆 Maturity"])
        
        with tab_main:
            if enable_sentiment:
                top_keys = [k for k,v in combined_counts.most_common(1000)]
                term_sentiments = get_sentiments(analyzer, tuple(top_keys))
                if proc_conf.compute_bigrams:
                     top_bg_keys = [" ".join(k) for k,v in scanner.global_bigrams.most_common(2000)]
                     term_sentiments.update(get_sentiments(analyzer, tuple(top_bg_keys)))
                c_color_func = create_sentiment_color_func(term_sentiments, pos_color, neg_color, neu_color, pos_threshold, neg_threshold)
                fig, _ = build_wordcloud_figure_from_counts(combined_counts, max_words, 800, 400, bg_color, colormap, combined_font_path, 42, c_color_func)
            else:
                term_sentiments = {}
                fig, _ = build_wordcloud_figure_from_counts(combined_counts, max_words, 800, 400, bg_color, colormap, combined_font_path, 42, None)
                
            st.pyplot(fig, use_container_width=True)
            st.download_button("📥 download combined png", fig_to_png_bytes(fig), "combined_wc.png", "image/png")
            
            c1, c2, c3, c4 = st.columns(4)
            
            c1.metric(
                "Total Tokens", 
                f"{text_stats['Total Tokens']:,}",
                help="The total count of all words processed after cleaning (removing stopwords, numbers, etc.). Represents the sheer volume of signal."
            )
            
            c2.metric(
                "Unique Vocab", 
                f"{text_stats['Unique Vocabulary']:,}",
                help="The count of distinct, unique words found. A higher number indicates a broader range of topics or more complex language."
            )
            
            c3.metric(
                "Docs/Rows", 
                f"{text_stats['Total Rows']:,}",
                help="The number of individual processing units (e.g., rows in a CSV, paragraphs in a PDF, or lines in a transcript)."
            )
            
            c4.metric(
                "Lexical Diversity", 
                f"{text_stats['Lexical Diversity']}",
                help="The Ratio of Unique Words to Total Words (Unique / Total). \n\n• High (>0.5): Dense information, varied vocabulary (e.g., Poetry, Abstracts).\n• Low (<0.1): Highly repetitive, consistent language (e.g., Logs, Legal Boilerplate)."
            )

        with tab_trend:
            if scanner.temporal_counts:
                st.markdown("#### Word Volume Over Time")
                trend_data = []
                for d_str, counts in scanner.temporal_counts.items():
                    trend_data.append({"Date": d_str, "Volume": sum(counts.values())})
                
                df_trend = pd.DataFrame(trend_data).sort_values("Date")
                st.line_chart(df_trend.set_index("Date"))
                
                st.markdown("#### Specific Term Trends")
                terms_to_plot = st.multiselect("Select terms to plot", [t for t, c in combined_counts.most_common(50)])
                if terms_to_plot:
                    term_trend_data = []
                    for d_str, counts in scanner.temporal_counts.items():
                        row = {"Date": d_str}
                        for t in terms_to_plot: row[t] = counts[t]
                        term_trend_data.append(row)
                    df_term_trend = pd.DataFrame(term_trend_data).sort_values("Date").set_index("Date")
                    st.line_chart(df_term_trend)
            else:
                st.info("No Date column was selected during scan (or no valid dates found).")

        with tab_ent:
            st.markdown("#### Top Entities (Polymorphic NER)")
            st.caption(
                "**What this finds:** Stakeholders (Who) and Systems (What). "
                "This engine adapts to capture **Standard Names** (e.g., 'John Doe'), "
                "**Acronyms** (e.g., 'DARPA'), and **Technical IDs** (e.g., 'COVID-19' or 'F-35')."
            )
            
            # smarter ratio filter
            # if "apple" (fruit) appears 100 times, and "Apple" (entity) appears 2 times, drop it
            # if "Apple" (Company) appears 100 times, keeps it
            refined_entities = Counter()
            for ent, count in scanner.entity_counts.items():
                lower_k = ent.lower()
                total_occurrences = scanner.global_counts.get(lower_k, 0)
                
                # if word appears in the text, check the ratio
                if total_occurrences > 0:
                    capitalization_ratio = count / total_occurrences
                    # keeps if it's capitalized >30% of the time OR it's a complex ID (digits/hyphens)
                    if capitalization_ratio > 0.3 or not ent.isalpha():
                        refined_entities[ent] = count
                else:
                    # if not in global_counts (e.g. stopped out), keeps it to be safe
                    refined_entities[ent] = count

            if refined_entities:
                ent_df = pd.DataFrame(refined_entities.most_common(50), columns=["Entity", "Count"])
                st.dataframe(ent_df, use_container_width=True)
                
                # simple entity cloud (safety wrapped)
                try:
                    fig_e, _ = build_wordcloud_figure_from_counts(refined_entities, 100, 800, 400, "#111111", "Pastel1", combined_font_path, 42, None)
                    st.pyplot(fig_e)
                except Exception as e:
                    st.warning(f"Could not generate Entity Cloud: {e}")
            else:
                st.info("No capitalized entities detected.")

        with tab_key:
            st.subheader("🔑 TF-IDF Keyphrases (The 'Technical DNA')")
            
            # 1 plain language explanation banner
            st.info(
                "**How to read this:** Unlike simple word counts, **TF-IDF** penalizes words that appear everywhere (like 'report' or 'email') and boosts words that are unique to specific incidents or documents. \n\n"
                "👉 **High Score** = Rare, specific, high-signal (e.g. 'Oximetry').\n"
                "👉 **Low Score** = Common, generic, low-signal.",
                icon="ℹ️"
            )
            
            df_tfidf = calculate_tfidf(scanner, 50)
            
            # 2 DataFrame with hover-over tooltips
            st.dataframe(
                df_tfidf, 
                use_container_width=True,
                column_config={
                    "Term": st.column_config.TextColumn("Term", help="The extracted vocabulary word."),
                    "TF (Count)": st.column_config.NumberColumn("TF (Count)", help="Term Frequency: Total number of times this word appears."),
                    "DF (Docs)": st.column_config.NumberColumn("DF (Docs)", help="Document Frequency: Number of distinct documents (or chunks) containing this word. Low DF = Specific."),
                    "Keyphrase Score": st.column_config.NumberColumn("Keyphrase Score", help="Mathematical Uniqueness. Higher = More 'Technical' and less 'Generic'.")
                }
            )

        with tab_mat:
            st.subheader("🏆 Maturity Assessment Engine")
            # 1. Initialize
            assessor = MaturityAssessor()
            model_options = assessor.get_model_names()
            # 2. Persona Selector
            c_sel1, c_sel2 = st.columns([1, 3])
            with c_sel1:
                selected_model = st.selectbox(
                    "Select Organization Persona:", model_options, index=0
                )
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
            with c_sel2:
                st.info(f"ℹ️ {assessor.get_model_desc(selected_model)}")
                st.caption(
                    "Maturity scoring uses the current cleaned token stream. "
                    "Settings such as Min Word Len, stopwords, lemmatization, "
                    "hyphen handling, and chat/HTML cleanup can change the result."
                )
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<
            # 3. Run Assessment
            maturity_result = assessor.assess(
                combined_counts, scanner.global_bigrams, selected_model
            )
            if maturity_result:
                result_type = maturity_result.get("type", "flat")
                if result_type == "flat":
                    # === ORIGINAL FLAT RENDERING (unchanged) ===
                    st.divider()
                    m_col1, m_col2 = st.columns([1, 2])
                    with m_col1:
                        score = maturity_result['overall_score']
                        dom_stage = maturity_result['dominant_stage']
                        gauge_color = dom_stage['color']
                        st.metric("Maturity Score (1.0 - 5.0)", f"{score} / 5.0")
                        st.markdown(
                            f"#### Phase: <span style='color:{gauge_color}'>"
                            f"{dom_stage['name']}</span>",
                            unsafe_allow_html=True
                        )
                        st.metric(
                            "Signal Density",
                            f"{maturity_result['total_signals_found']} words",
                            help="Count of relevant vocabulary words found."
                        )
                    with m_col2:
                        fig_mat = assessor.render_radar_chart(maturity_result)
                        if fig_mat:
                            st.pyplot(fig_mat, use_container_width=True)
                        st.caption(
                            "**Analyst Tip:** A balanced polygon (round) indicates "
                            "alignment between Strategy and Execution. A sharp "
                            "'spike' indicates a capability gap in adjacent areas."
                        )
                    st.subheader("🔍 Linguistic Drivers")
                    cols = st.columns(5)
                    levels_ref = maturity_result['levels_ref']
                    for i in range(1, 6):
                        with cols[i - 1]:
                            lvl_data = levels_ref[i]
                            st.markdown(f"**L{i}: {lvl_data['name']}**")
                            terms = lvl_data['terms']
                            found = {
                                t: combined_counts[t]
                                for t in terms if t in combined_counts
                            }
                            if found:
                                top_words = sorted(
                                    found.items(), key=lambda x: x[1], reverse=True
                                )[:4]
                                for w, c in top_words:
                                    st.caption(f"{w} ({c})")
                            else:
                                st.caption("(No signals)")
                elif result_type == "domain_based":
                    # === NEW DOMAIN-BASED RENDERING ===
                    st.divider()
                    # --- Top-level metrics ---
                    m_col1, m_col2, m_col3 = st.columns(3)
                    with m_col1:
                        st.metric(
                            "Composite Maturity (1.0 - 3.0)",
                            f"{maturity_result['overall_score']} / 3.0"
                        )
                    with m_col2:
                        st.metric(
                            "Domains Assessed",
                            f"{maturity_result['domains_assessed']} / "
                            f"{maturity_result['domains_total']}"
                        )
                    with m_col3:
                        st.metric(
                            "Total Signals",
                            f"{maturity_result['total_signals_found']}"
                        )
                    # --- Radar Chart ---
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
                    # --- Radar Chart ---
                    st.markdown("#### 🕸️ Domain Maturity Radar")

                    with st.expander("🎨 Radar font controls", expanded=False):
                        font_col1, font_col2, font_col3 = st.columns(3)

                        with font_col1:
                            radar_font_family = st.selectbox(
                                "Font family",
                                ["sans-serif", "serif", "monospace", "DejaVu Sans", "Arial"],
                                index=0,
                                key="radar_font_family",
                            )
                            radar_label_size = st.slider(
                                "Domain label size",
                                min_value=6,
                                max_value=16,
                                value=9,
                                step=1,
                                key="radar_label_size",
                            )

                        with font_col2:
                            radar_label_color = st.color_picker(
                                "Label color",
                                value="#222222",
                                key="radar_label_color",
                            )
                            radar_wrap_width = st.slider(
                                "Domain label wrap width",
                                min_value=6,
                                max_value=24,
                                value=14,
                                step=1,
                                key="radar_wrap_width",
                            )

                        with font_col3:
                            radar_show_tier_labels = st.checkbox(
                                "Show tier names",
                                value=True,
                                key="radar_show_tier_labels",
                            )
                            radar_tier_label_angle = st.slider(
                                "Tier label angle",
                                min_value=0,
                                max_value=359,
                                value=35,
                                step=5,
                                key="radar_tier_label_angle",
                            )

                    fig_radar = assessor.render_domain_radar_chart(
                        maturity_result,
                        font_family=radar_font_family,
                        label_size=radar_label_size,
                        label_color=radar_label_color,
                        wrap_width=radar_wrap_width,
                        show_tier_labels=radar_show_tier_labels,
                        tier_label_angle=radar_tier_label_angle,
                    )

                    if fig_radar:
                        st.pyplot(fig_radar, use_container_width=True)
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<
                    st.caption(
                        "🔴 Foundational (<1.5) | 🟠 Advanced (1.5–2.5) | "
                        "🟢 Leading Edge (>2.5)"
                    )
                    # --- Stacked Bar Breakdown ---
                    st.markdown("#### 📊 Tier Distribution by Domain")
                    fig_bar = assessor.render_domain_breakdown_chart(maturity_result)
                    if fig_bar:
                        st.pyplot(fig_bar, use_container_width=True)
                    # --- Per-Domain Detail Cards ---
#>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>
                    st.markdown("#### 🔍 Domain Details")
                    domain_results = maturity_result["domain_results"]
                    for domain_key in sorted(domain_results.keys()):
                        dr = domain_results[domain_key]

                        if dr["signals"] == 0:
                            indicator = "⚪"
                        elif dr["score"] >= 2.5:
                            indicator = "🟢"
                        elif dr["score"] >= 1.5:
                            indicator = "🟠"
                        else:
                            indicator = "🔴"

                        with st.expander(
                            f"{indicator} {dr['name']} — "
                            f"{dr['tier_label']} ({dr['score']}/3.0, "
                            f"{dr['signals']} signals)"
                        ):
                            if dr["signals"] == 0:
                                st.info(
                                    "No matching vocabulary from this dataset was mapped to this domain yet. "
                                    "This usually means either the source material did not discuss this area, "
                                    "or the current domain vocabulary is too narrow for the way it was discussed."
                                )
                                continue

                            dist = dr["distribution"]
                            dc1, dc2, dc3 = st.columns(3)
                            dc1.metric("Foundational", f"{dist.get(1, 0):.0%}")
                            dc2.metric("Advanced", f"{dist.get(2, 0):.0%}")
                            dc3.metric("Leading Edge", f"{dist.get(3, 0):.0%}")
                            st.markdown("**Top Linguistic Drivers:**")
                            tier_names = {
                                1: "Foundational", 2: "Advanced", 3: "Leading Edge"
                            }
                            for tier_num in [1, 2, 3]:
                                drivers = dr["drivers"].get(tier_num, [])
                                if drivers:
                                    driver_str = ", ".join(
                                        [f"`{t}` ({c})" for t, c in drivers]
                                    )
                                    st.caption(
                                        f"**{tier_names[tier_num]}:** {driver_str}"
                                    )
#<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<<
                    # ============================================
                    # PHASE 10: DOWNLOAD ASSESSMENT (inline here)
                    # ============================================
                    st.divider()
                    st.markdown("#### 💾 Export Assessment")
                    client_name_input = st.text_input(
                        "Client Name (for export filename):",
                        value="Client",
                        key="mat_client_name"
                    )
                    # Sanitize client name for filename
                    safe_client = re.sub(r'[^\w\s-]', '', client_name_input).strip()
                    safe_client = re.sub(r'[\s]+', '_', safe_client)
                    export_date = datetime.now().strftime("%Y-%m-%d")
                    export_filename = f"{safe_client}_Maturity_Assessment_{export_date}.json"
                    # Build export payload
                    export_payload = {
                        "schema_version": "1.0",
                        "client_name": client_name_input.strip(),
                        "assessment_date": export_date,
                        "model_name": selected_model,
                        "composite_score": maturity_result["overall_score"],
                        "max_score": maturity_result["max_score"],
                        "domains_assessed": maturity_result["domains_assessed"],
                        "domains_total": maturity_result["domains_total"],
                        "total_signals": maturity_result["total_signals_found"],
                        "domain_results": {}
                    }
                    for dk, dr in maturity_result["domain_results"].items():
                        export_payload["domain_results"][dk] = {
                            "name": dr["name"],
                            "short": dr["short"],
                            "score": dr["score"],
                            "tier_label": dr["tier_label"],
                            "distribution": {
                                "foundational": round(dr["distribution"].get(1, 0), 4),
                                "advanced": round(dr["distribution"].get(2, 0), 4),
                                "leading_edge": round(dr["distribution"].get(3, 0), 4)
                            },
                            "signals": dr["signals"],
                            "top_drivers": {
                                "foundational": [
                                    f"{t} ({c})" for t, c in dr["drivers"].get(1, [])
                                ],
                                "advanced": [
                                    f"{t} ({c})" for t, c in dr["drivers"].get(2, [])
                                ],
                                "leading_edge": [
                                    f"{t} ({c})" for t, c in dr["drivers"].get(3, [])
                                ]
                            }
                        }
                    export_json_str = json.dumps(export_payload, indent=2)
                    st.download_button(
                        label=f"⬇️ Download: {export_filename}",
                        data=export_json_str,
                        file_name=export_filename,
                        mime="application/json"
                    )
                    # =====================================================
                    # PHASE 11: UPLOAD SNAPSHOTS + LONGITUDINAL TRACKING
                    # =====================================================
                    st.divider()
                    st.markdown("#### 📈 Longitudinal Progress Tracking")
                    st.caption(
                        "Upload previous assessment JSON files to track maturity "
                        "progress over time. The current live assessment (if running) "
                        "will be included automatically."
                    )
                    uploaded_snapshots = st.file_uploader(
                        "Upload previous assessment snapshots (.json):",
                        type=["json"],
                        accept_multiple_files=True,
                        key="mat_snapshot_upload"
                    )
                    # Parse uploaded snapshots
                    all_snapshots = []
                    parse_errors = []
                    if uploaded_snapshots:
                        for uf in uploaded_snapshots:
                            try:
                                content = json.loads(uf.read())
                                uf.seek(0)  # reset for potential re-read
                                # Validate schema
                                if "schema_version" not in content:
                                    parse_errors.append(
                                        f"⚠️ {uf.name}: Missing schema_version. Skipping."
                                    )
                                    continue
                                if "assessment_date" not in content:
                                    parse_errors.append(
                                        f"⚠️ {uf.name}: Missing assessment_date. Skipping."
                                    )
                                    continue
                                all_snapshots.append(content)
                            except json.JSONDecodeError:
                                parse_errors.append(
                                    f"⚠️ {uf.name}: Invalid JSON. Skipping."
                                )
                    # Append current live assessment as latest snapshot
                    if maturity_result and maturity_result.get("type") == "domain_based":
                        live_snapshot = {
                            "schema_version": "1.0",
                            "client_name": client_name_input.strip(),
                            "assessment_date": export_date,
                            "composite_score": maturity_result["overall_score"],
                            "domain_results": {}
                        }
                        for dk, dr in maturity_result["domain_results"].items():
                            live_snapshot["domain_results"][dk] = {
                                "name": dr["name"],
                                "short": dr["short"],
                                "score": dr["score"],
                                "tier_label": dr["tier_label"],
                                "signals": dr["signals"]
                            }
                        all_snapshots.append(live_snapshot)
                    # Sort by date
                    all_snapshots.sort(
                        key=lambda s: s.get("assessment_date", "0000-00-00")
                    )
                    # Show parse errors
                    for err in parse_errors:
                        st.warning(err)
                    if len(all_snapshots) >= 2:
                        st.success(
                            f"📊 {len(all_snapshots)} snapshots loaded "
                            f"(including current assessment)."
                        )
                        # --- Composite Trend ---
                        st.markdown("##### Composite Score Trend")
                        fig_trend = assessor.render_trend_chart(all_snapshots)
                        if fig_trend:
                            st.pyplot(fig_trend, use_container_width=True)
                        # --- Per-Domain Trend ---
                        st.markdown("##### Per-Domain Score Trend")
                        fig_domain_trend = assessor.render_domain_trend_chart(
                            all_snapshots
                        )
                        if fig_domain_trend:
                            st.pyplot(fig_domain_trend, use_container_width=True)
                        # --- Domain Movement Table ---
                        st.markdown("##### Domain Movement (First → Latest)")
                        first = all_snapshots[0]
                        latest = all_snapshots[-1]
                        movement_data = []
                        all_domain_keys = set(
                            list(first.get("domain_results", {}).keys()) +
                            list(latest.get("domain_results", {}).keys())
                        )
                        for dk in sorted(all_domain_keys):
                            first_dr = first.get("domain_results", {}).get(dk, {})
                            latest_dr = latest.get("domain_results", {}).get(dk, {})
                            first_score = first_dr.get("score", 0.0)
                            latest_score = latest_dr.get("score", 0.0)
                            delta = round(latest_score - first_score, 2)
                            if delta > 0:
                                direction = "⬆️ Improved"
                            elif delta < 0:
                                direction = "⬇️ Declined"
                            else:
                                direction = "➡️ No Change"
                            short_name = (
                                latest_dr.get("short") or
                                first_dr.get("short") or dk
                            )
                            movement_data.append({
                                "Domain": short_name,
                                "First Score": first_score,
                                "Latest Score": latest_score,
                                "Change": delta,
                                "Direction": direction
                            })
                        if movement_data:
                            st.dataframe(
                                movement_data,
                                use_container_width=True,
                                hide_index=True
                            )
                    elif len(all_snapshots) == 1:
                        st.info(
                            "Upload at least one previous snapshot to see "
                            "progress trends. Currently showing only the "
                            "live assessment."
                        )
                    else:
                        st.info(
                            "No snapshots loaded. Run an assessment and export "
                            "it, then upload past exports here to track progress."
                        )
            else:
                st.warning(
                    "No sufficient vocabulary found for this specific Maturity "
                    "Model. Try switching Personas or adding more text data."
                )

        st.divider()
        
        # advanced sections
        
        if enable_sentiment and beta_dist:
            st.subheader("⚖️ Bayesian Sentiment Inference")
            with st.expander("🧠 How to read this chart (and why it matters)", expanded=False):
                st.markdown("""
                **The Problem:** Standard sentiment analysis gives you a single number (e.g., "52% Positive"). But is that 52% based on 5 tweets or 5 million? A single number hides that uncertainty.
                
                **The Solution:** This chart calculates the **Probability** of the true sentiment.
                *   **The Curve (PDF):** Represents likelihood. The higher the peak, the more likely that specific sentiment score is the "truth."
                *   **The Shape:** 
                    *   **Narrow & Tall:** We have lots of data. We are highly confident the sentiment is exactly here.
                    *   **Wide & Flat:** We don't have enough data. The true sentiment could be almost anything.
                *   **The Green Zone (95% CI):** There is a 95% probability the "True" sentiment falls within this range. 
                
                **Decision Tip:** If the green zone is very wide (e.g., spanning 30% to 70%), **do not** make business decisions based on sentiment yet; you need more data.
                """)

            bayes_result = perform_bayesian_sentiment_analysis(combined_counts, term_sentiments, pos_threshold, neg_threshold)
            if bayes_result:
                b_col1, b_col2 = st.columns([1, 2])
                with b_col1:
                    st.metric("Positive Words Observed", f"{bayes_result['pos_count']:,}")
                    st.metric("Negative Words Observed", f"{bayes_result['neg_count']:,}")
                    st.info(f"Mean Expected Positive Rate: **{bayes_result['mean_prob']:.1%}**")
                    st.success(f"95% Credible Interval:\n**{bayes_result['ci_low']:.1%} — {bayes_result['ci_high']:.1%}**")
                with b_col2:
                    fig_bayes, ax_bayes = plt.subplots(figsize=(8, 4))
                    ax_bayes.plot(bayes_result['x_axis'], bayes_result['pdf_y'], lw=2, color='blue', label='Posterior PDF')
                    ax_bayes.fill_between(bayes_result['x_axis'], 0, bayes_result['pdf_y'], 
                                        where=(bayes_result['x_axis'] > bayes_result['ci_low']) & (bayes_result['x_axis'] < bayes_result['ci_high']),
                                        color='green', alpha=0.3, label='95% Credible Interval')
                    ax_bayes.set_title("Bayesian Update of Sentiment Confidence", fontsize=10)
                    ax_bayes.legend()
                    ax_bayes.grid(True, alpha=0.2)
                    st.pyplot(fig_bayes)
                    plt.close(fig_bayes)

        show_graph = proc_conf.compute_bigrams and scanner.global_bigrams and st.checkbox("🕸️ Show Network Graph & Advanced Analytics", value=True)
        if show_graph:
            st.subheader("🔗 Network Graph")
            with st.expander("🛠️ Graph Settings & Physics", expanded=False):
                c1, c2, c3 = st.columns(3)
                min_edge_weight = c1.slider(
                    "Min Link Frequency", 2, 100, 2, 
                    help="Minimum shared occurrences required to draw a line. Increase this to remove weak connections and 'de-clutter' the graph."
                )
                max_nodes_graph = c1.slider(
                    "Max Nodes", 10, 200, 80, 
                    help="Hard limit on the number of words displayed. Lower this to focus only on the absolute most vital connections."
                )
                repulsion_val = c2.slider(
                    "Repulsion", 100, 3000, 1000, 
                    help="Physics Force: How strongly nodes push away from each other. Increase this if the graph looks like a tight ball."
                )
                edge_len_val = c2.slider(
                    "Edge Length", 50, 500, 250, 
                    help="Target length for the connecting lines. Increase this to space out distinct clusters."
                )
                physics_enabled = c3.checkbox(
                    "Enable Physics", True, 
                    help="If checked, the graph simulates gravity to organize itself. Uncheck to freeze the nodes in place."
                )
                directed_graph = c3.checkbox(
                    "Directed Arrows", False, 
                    help="Draws arrows (->) to show word order/flow, rather than just simple connections."
                )
                color_options = ["Community (Topic)", "Sentiment"]
                # Add maturity domain coloring if available
                if ('maturity_result' in locals() and maturity_result
                    and maturity_result.get("type") == "domain_based"):
                    color_options.append("Maturity Domain")
                color_mode = c3.radio(
                    "Color By:", color_options, index=0,
                    help="Community: Colors by structural cluster.\n"
                    "Sentiment: Colors by positive/negative.\n"
                    "Maturity Domain: Colors by which maturity domain the word belongs to (if 12-domain model is active)."
                )

            G = nx.DiGraph() if directed_graph else nx.Graph()
            filtered_bigrams = {k: v for k, v in scanner.global_bigrams.items() if v >= min_edge_weight}
            sorted_connections = sorted(filtered_bigrams.items(), key=lambda x: x[1], reverse=True)[:max_nodes_graph]
            
            if sorted_connections:
                G.add_edges_from((src, tgt, {'weight': w}) for (src, tgt), w in sorted_connections)
                try: deg_centrality = nx.degree_centrality(G)
                except: deg_centrality = {n: 1 for n in G.nodes()}
                community_map = {}
                ai_cluster_info = ""
                
                if color_mode == "Community (Topic)":
                    G_undir = G.to_undirected() if directed_graph else G
                    try:
                        communities = nx_comm.greedy_modularity_communities(G_undir)
                        cluster_descriptions = []
                        for group_id, community in enumerate(communities):
                            top_in_cluster = sorted(list(community), key=lambda x: combined_counts[x], reverse=True)[:5]
                            cluster_descriptions.append(f"- Cluster {group_id+1}: {', '.join(top_in_cluster)}")
                            for node in community: community_map[node] = group_id
                        ai_cluster_info = "\n".join(cluster_descriptions)
                    except: pass

                community_colors = ["#FF4B4B", "#4589ff", "#ffa421", "#3cdb82", "#8b46ff", "#ff4b9f", "#00c0f2"]
                nodes, edges = [], []
                for node_id in G.nodes():
                    size = 15 + (deg_centrality.get(node_id, 0) * 80)
                    node_color = "#808080"
                    if color_mode == "Sentiment" and enable_sentiment:
                        s = term_sentiments.get(node_id, 0)
                        if s >= pos_threshold: node_color = pos_color
                        elif s <= neg_threshold: node_color = neg_color
                    elif color_mode == "Community (Topic)":
                        gid = community_map.get(node_id, 0)
                        node_color = community_colors[gid % len(community_colors)]
                    elif color_mode == "Maturity Domain" and maturity_result:
                        # Approximation: a word may appear in multiple domains, so the first mapped domain wins.
                        domain_colors_map = {
                            "01": "#e6194b", "02": "#3cb44b", "03": "#ffe119",
                            "04": "#4363d8", "05": "#f58231", "06": "#911eb4",
                            "07": "#42d4f4", "08": "#f032e6", "09": "#bfef45",
                            "10": "#fabed4", "11": "#469990", "12": "#dcbeff"
                        }
                        term_domain_map = {}
                        for dk, dd in maturity_result.get("domains_ref", {}).items():
                            prefix = dk[:2]
                            for tier_num, tier_data in dd["tiers"].items():
                                for term in tier_data.get("terms", set()):
                                    if term not in term_domain_map:
                                        term_domain_map[term] = prefix
                                for phrase in tier_data.get("phrases", []):
                                    for word in phrase.split():
                                        if word not in term_domain_map:
                                            term_domain_map[word] = prefix
                        matched_domain = term_domain_map.get(node_id)
                        if matched_domain:
                            node_color = domain_colors_map.get(matched_domain, "#808080")
                        else:
                            node_color = "#808080"

                    # re-added font config for white, legible text
                    nodes.append(Node(
                        id=node_id, 
                        label=node_id, 
                        size=size, 
                        color=node_color,
                        font={'color': 'white', 'size': 20, 'strokeWidth': 2, 'strokeColor': '#000000'}
                    ))

                for (source, target), weight in sorted_connections:
                    width = 1 + math.log(weight) * 0.8
                    edges.append(Edge(source=source, target=target, width=width, color="#e0e0e0"))
                
                # re-added interaction dict for zoom/pan buttons
                config = Config(
                    width=1000, 
                    height=700, 
                    directed=directed_graph, 
                    physics=physics_enabled, 
                    hierarchy=False, 
                    interaction={"navigationButtons": True, "zoomView": True}, 
                    physicsSettings={"solver": "forceAtlas2Based", "forceAtlas2Based": {"gravitationalConstant": -abs(repulsion_val), "springLength": edge_len_val, "springConstant": 0.05, "damping": 0.4}}
                )
                
                st.info("💡 **Navigation Tip:** Use the buttons in the **bottom-right** of the graph to Zoom & Pan.")
                agraph(nodes=nodes, edges=edges, config=config)

                # [NEW] Gephi Export
                if st.button("📥 Download Graph File (.gexf)"):
                    try:
                        # Write to memory buffer
                        gexf_buffer = io.BytesIO()
                        nx.write_gexf(G, gexf_buffer)
                        gexf_buffer.seek(0)
                        st.download_button("Click to Save GEXF", gexf_buffer, "network.gexf", "application/xml")
                    except Exception as e:
                        st.error(f"Export failed: {e}")
                
                # graph analytics tabs
                tab_g1, tab_g2, tab_g3, tab_g4 = st.tabs(["Basic Stats", "Top Nodes", "Text Stats", "🔥 Heatmap"])
                with tab_g1:
                    col_b1, col_b2, col_b3 = st.columns(3)
                    col_b1.metric(
                        "Nodes", 
                        G.number_of_nodes(),
                        help="The number of unique concepts (dots) currently displayed in the network."
                    )
                    col_b2.metric(
                        "Edges", 
                        G.number_of_edges(),
                        help="The number of connections (lines) between concepts. A connection indicates these words appear together frequently."
                    )
                    try: 
                        col_b3.metric(
                            "Density", 
                            f"{nx.density(G):.4f}",
                            help="A score from 0 to 1 indicating how interconnected the network is.\n\n• Low (<0.1): Distinct, separate topics.\n• High (>0.5): Everything is related to everything (a 'hairball')."
                        )
                    except: pass
                with tab_g2:
                    node_weights = {n: 0 for n in G.nodes()}
                    for u, v, data in G.edges(data=True):
                        w = data.get('weight', 1)
                        node_weights[u] += w
                        node_weights[v] += w
                    st.dataframe(pd.DataFrame(list(node_weights.items()), columns=["Node", "Weighted Degree"]).sort_values("Weighted Degree", ascending=False).head(50), use_container_width=True)
                with tab_g3:
                     c1, c2, c3 = st.columns(3)
                     c1.metric(
                        "Total Tokens", 
                        f"{text_stats['Total Tokens']:,}",
                        help="The total count of all words processed after cleaning (removing stopwords, numbers, etc.). Represents the sheer volume of signal."
                     )
                     c2.metric(
                        "Unique Vocab", 
                        f"{text_stats['Unique Vocabulary']:,}",
                        help="The count of distinct, unique words found. A higher number indicates a broader range of topics or more complex language."
                     )
                     c3.metric(
                        "Lexical Diversity", 
                        f"{text_stats['Lexical Diversity']}",
                        help="The Ratio of Unique Words to Total Words (Unique / Total). \n\n• High (>0.5): Dense information, varied vocabulary (e.g., Poetry, Abstracts).\n• Low (<0.1): Highly repetitive, consistent language (e.g., Logs, Legal Boilerplate)."
                     )
                with tab_g4:
                    # hybrid heatmap-QR generator
                    viz_mode = st.radio("Visualization Mode", ["Standard Heatmap", "Hybrid Signature (Scanable)"], horizontal=True, label_visibility="collapsed")
                    
                    top_20 = [w for w, c in combined_counts.most_common(20)]
                    
                    if len(top_20) > 1:
                        # generating matrix data
                        mat = np.zeros((len(top_20), len(top_20)))
                        for i, w1 in enumerate(top_20):
                            for j, w2 in enumerate(top_20):
                                if i != j: mat[i][j] = scanner.global_bigrams.get((w1, w2), 0) + scanner.global_bigrams.get((w2, w1), 0)
                        
                        # to plot heatmap to a PIL image (memory buffer)
                        fig_h, ax_h = plt.subplots(figsize=(10, 10)) 
                        ax_h.imshow(mat, cmap=colormap, interpolation='nearest') 
                        
                        if viz_mode == "Hybrid Signature (Scanable)":
                            ax_h.axis('off')
                            plt.tight_layout(pad=0)
                        else:
                            ax_h.set_xticks(np.arange(len(top_20)))
                            ax_h.set_yticks(np.arange(len(top_20)))
                            ax_h.set_xticklabels(top_20, rotation=45, ha="right")
                            ax_h.set_yticklabels(top_20)
                        
                        buf = BytesIO()
                        fig_h.savefig(buf, format="png", bbox_inches="tight", pad_inches=0)
                        buf.seek(0)
                        
                        if viz_mode == "Standard Heatmap":
                            st.pyplot(fig_h)
                        
                        elif viz_mode == "Hybrid Signature (Scanable)":
                            if qrcode is None:
                                st.error("🚨 Please install: `pip install qrcode[pil]`")
                            else:
                                from PIL import Image, ImageEnhance
                                
                                # to allow for custom "steganographic" payload (with char limit)
                                custom_msg = st.text_input(
                                    "🔒 Encode Custom Payload (Optional)", 
                                    placeholder="Leave empty for standard metadata...",
                                    max_chars=1000,
                                    help="Limit: ~1000 characters to ensure QR readability with high error correction."
                                )
                                
                                # prepping heatmap
                                heatmap_img = Image.open(buf).convert("RGBA")
                                enhancer = ImageEnhance.Brightness(heatmap_img)
                                heatmap_img = enhancer.enhance(1.5) 
                                
                                # defining payload
                                if custom_msg.strip():
                                    signature_payload = custom_msg
                                else:
                                    signature_payload = (
                                        f"SIGNAL FOUNDRY\nRef: {st.session_state.get('last_sketch_hash', 'SESSION')}\n"
                                        f"Top: {', '.join(top_20[:3])}"
                                    )
                                
                                # to generate QR with safety check
                                try:
                                    qr = qrcode.QRCode(error_correction=qrcode.constants.ERROR_CORRECT_H, border=1)
                                    qr.add_data(signature_payload)
                                    qr.make(fit=True)
                                    qr_img = qr.make_image(fill_color="black", back_color="transparent").convert("RGBA")
                                    
                                    # composite
                                    heatmap_resized = heatmap_img.resize(qr_img.size)
                                    hybrid_img = Image.alpha_composite(heatmap_resized, qr_img)
                                    
                                    c1, c2 = st.columns([2, 1])
                                    with c1:
                                        st.image(hybrid_img, caption="Scan to verify data (or read payload).", use_container_width=True)
                                    with c2:
                                        st.markdown("### 🧬 Hybrid Signature")
                                        if custom_msg.strip():
                                            st.warning("⚠️ **Mode:** Custom Payload Active")
                                        else:
                                            st.info("ℹ️ **Mode:** Standard Metadata")
                                        
                                        st.caption("The colors represent the data relationships (Heatmap). The dark overlay pattern encodes the message (QR).")
                                        
                                        final_buf = BytesIO()
                                        hybrid_img.save(final_buf, format="PNG")
                                        st.download_button("📥 Download Signature", final_buf.getvalue(), "heatmap_signature.png", "image/png")
                                        
                                except Exception as e:
                                    st.error("❌ **Payload Too Large:** The text is too long to fit into a High-Security QR code. Please shorten your message.")
                    else:
                        st.info("Not enough data to generate signature.")
                    #

        st.subheader("🔍 Bayesian Theme Discovery")
        if len(scanner.topic_docs) > 5 and DictVectorizer:
            with st.spinner(f"Running {topic_model_type} Topic Modeling..."):
                topics = perform_topic_modeling(scanner.topic_docs, n_topics, topic_model_type)
            if topics:
                cols = st.columns(len(topics))
                for idx, topic in enumerate(topics):
                    with cols[idx]:
                        st.markdown(f"**Topic {topic['id']}**")
                        for w in topic['words']: st.markdown(f"`{w}`")
        else:
            st.info("Needs more data/docs to model topics.")

        # detailed frequency tables
        st.divider()
        st.subheader(f"📊 Frequency Tables (Top {top_n})")
        most_common = combined_counts.most_common(top_n)
        data = []
        if enable_sentiment:
            for w, f in most_common:
                score = term_sentiments.get(w, 0.0)
                category = get_sentiment_category(score, pos_threshold, neg_threshold)
                data.append([w, f, score, category])
        else:
            data = [[w, f] for w, f in most_common]

        cols = ["word", "count"] + (["sentiment", "category"] if enable_sentiment else [])
        st.dataframe(pd.DataFrame(data, columns=cols), use_container_width=True)
        
        if proc_conf.compute_bigrams and scanner.global_bigrams:
            st.write("Bigrams (By Frequency)")
            top_bg = scanner.global_bigrams.most_common(top_n)
            bg_data = []
            if enable_sentiment:
                for bg_tuple, f in top_bg:
                    bg_str = " ".join(bg_tuple)
                    score = term_sentiments.get(bg_str, 0.0)
                    category = get_sentiment_category(score, pos_threshold, neg_threshold)
                    bg_data.append([bg_str, f, score, category])
            else:
                bg_data = [[" ".join(bg), f] for bg, f in top_bg]
            bg_cols = ["bigram", "count"] + (["sentiment", "category"] if enable_sentiment else [])
            st.dataframe(pd.DataFrame(bg_data, columns=bg_cols), use_container_width=True)

            # NPMI in expander (original style)
            with st.expander("🔬 Phrase Significance (NPMI Score)", expanded=False):
                st.markdown("""
                **NPMI (Normalized Pointwise Mutual Information)** finds words that *belong* together, rather than just words that appear often.
                *   High Score (> 0.5): Strong association (e.g., "Artificial Intelligence").
                *   Low Score (< 0.1): Random association (e.g., "of the").
                """)
                df_npmi = calculate_npmi(scanner.global_bigrams, combined_counts, scanner.total_rows_processed)
                st.dataframe(df_npmi.head(top_n), use_container_width=True)

    # --- AI analyst (restored full mode)
    if combined_counts and st.session_state['authenticated']:
        st.divider()
        st.subheader("🤖 AI Analyst")
        st.caption(
            "**Privacy Note:** This AI does *not* read your raw documents. "
            "It interprets the statistical 'Sketch' (frequencies, connections, and clusters) generated by this engine. "
            "This allows for high-level pattern recognition without exposing the full text content."
        )
        
        top_u = [w for w, c in combined_counts.most_common(50)]
        top_b = [" ".join(bg) for bg, c in scanner.global_bigrams.most_common(20)]
        # Build AI context string
        ai_ctx_str = f"Top Words: {', '.join(top_u)}\nTop Bigrams: {', '.join(top_b)}\nGraph Clusters: {locals().get('ai_cluster_info', 'N/A')}"
        # Enrich with maturity data if available
        if 'maturity_result' in locals() and maturity_result:
            if maturity_result.get("type") == "domain_based":
                domain_summary_parts = []
                for dk in sorted(maturity_result["domain_results"].keys()):
                    dr = maturity_result["domain_results"][dk]
                    if dr["signals"] > 0:
                        domain_summary_parts.append(
                            f"{dr['short']}: {dr['score']}/3.0 ({dr['tier_label']})"
                        )
                if domain_summary_parts:
                    ai_ctx_str += f"\n\n12-Domain Maturity Scores:\n" + "\n".join(domain_summary_parts)
                    ai_ctx_str += f"\nComposite Score: {maturity_result['overall_score']}/3.0"
            else:
                ai_ctx_str += f"\nMaturity Score: {maturity_result.get('overall_score', 'N/A')}/5.0"
        
        col_ai_1, col_ai_2 = st.columns(2)
        
        with col_ai_1:
            st.markdown("**1. One-Click Theme Detection**")
            if st.button("✨ Identify Key Themes", type="primary"):
                with st.status("Analyzing..."):
                    system_prompt = (
                        "You are a qualitative data analyst. Analyze the provided word frequency lists to identify "
                        "3 key themes, potential anomalies, and a summary of the subject matter. "
                        "If 12-Domain Maturity Scores are provided, also identify the strongest and weakest domains, "
                        "and suggest 2-3 coaching priorities based on the gaps."
                    )
                    user_prompt = f"Data Context:\n{ai_ctx_str}"
                    response = call_llm_and_track_cost(system_prompt, user_prompt, ai_config)
                    st.session_state["ai_response"] = response
                    st.rerun()

        with col_ai_2:
            st.markdown("**2. Ask the Data**")
            user_question = st.text_area("Ask a specific question:", height=100, placeholder="e.g., 'What are the main complaints about pricing?'")
            if st.button("Ask Question"):
                if user_question.strip():
                    with st.status("Thinking..."):
                        system_prompt = "You are an expert analyst. Answer the user's question based ONLY on the provided summary statistics (word counts and associations). If you cannot answer from the data, say so."
                        user_prompt = f"Data Context:\n{ai_ctx_str}\n\nUser Question: {user_question}"
                        response = call_llm_and_track_cost(system_prompt, user_prompt, ai_config)
                        st.session_state["ai_response"] = f"**Q: {user_question}**\n\n{response}"
                        st.rerun()
                else:
                    st.warning("Please enter a question.")

        if st.session_state["ai_response"]:
            st.divider()
            st.markdown("### 📋 AI Output")
            st.markdown(st.session_state["ai_response"])

st.markdown("---")
st.markdown(
    "<div style='text-align: center; color: #808080; font-size: 12px;'>"
    "Open Source software licensed under the MIT License."
    "</div>", 
    unsafe_allow_html=True
)
